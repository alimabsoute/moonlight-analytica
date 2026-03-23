#!/usr/bin/env python3
"""
Tastyr IQ — Investor-Grade Outreach Strategy Deck Builder
Generates a clean, modern 16:9 PowerPoint deck using python-pptx.

Output: tastyr-iq-outreach-deck.pptx
Design: White backgrounds, Calibri, green/amber accents, metric blocks
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Globals ──────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF7, 0xF9, 0xFC)
CHARCOAL     = RGBColor(0x1F, 0x29, 0x37)
GREEN        = RGBColor(0x16, 0xA3, 0x4A)
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_GRAY   = RGBColor(0xE5, 0xE7, 0xEB)
MID_GRAY     = RGBColor(0x9C, 0xA3, 0xAF)
SOFT_RED     = RGBColor(0xEF, 0x44, 0x44)
DARK_GREEN   = RGBColor(0x15, 0x80, 0x3D)

FONT = "Calibri"

# ── Helpers ──────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color=WHITE):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=CHARCOAL, alignment=PP_ALIGN.LEFT,
                font_name=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """Add a textbox with a single run of formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_rich_textbox(slide, left, top, width, height, runs,
                     alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                     line_spacing=None):
    """Add a textbox with multiple runs (list of dicts with text, size, bold, color)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing is not None:
        p.line_spacing = line_spacing
    for r in runs:
        run = p.add_run()
        run.text = r.get("text", "")
        run.font.size = Pt(r.get("size", 14))
        run.font.bold = r.get("bold", False)
        run.font.color.rgb = r.get("color", CHARCOAL)
        run.font.name = r.get("font", FONT)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                          line_spacing=None):
    """
    lines: list of dicts, each representing a paragraph.
    Each dict has: text, size, bold, color, alignment (optional).
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = line.get("alignment", alignment)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if "space_before" in line:
            p.space_before = line["space_before"]
        if "space_after" in line:
            p.space_after = line["space_after"]
        run = p.add_run()
        run.text = line.get("text", "")
        run.font.size = Pt(line.get("size", 14))
        run.font.bold = line.get("bold", False)
        run.font.color.rgb = line.get("color", CHARCOAL)
        run.font.name = line.get("font", FONT)
    return txBox


def add_accent_line(slide, left, top, width, color=GREEN, height=Pt(3)):
    """Add a thin horizontal accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_metric_block(slide, left, top, width, height, number, label,
                     number_color=GREEN, label_color=MID_GRAY,
                     number_size=36, label_size=12, bg_color=OFF_WHITE):
    """Add a rounded-rect metric block with large number and small label."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = LIGHT_GRAY
    shape.line.width = Pt(0.5)
    # Reduce corner rounding
    shape.adjustments[0] = 0.05

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    # Number paragraph
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(number_size)
    run.font.bold = True
    run.font.color.rgb = number_color
    run.font.name = FONT

    # Label paragraph
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(label_size)
    run2.font.bold = False
    run2.font.color.rgb = label_color
    run2.font.name = FONT

    return shape


def add_step_circle(slide, left, top, size, number, color=GREEN):
    """Add a numbered circle for process flows."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT
    return shape


def add_arrow_right(slide, left, top, width=Inches(0.5), height=Inches(0.3)):
    """Add a right-pointing arrow connector."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total=21):
    """Add a small slide number in the bottom-right corner."""
    add_textbox(
        slide,
        Inches(12.2), Inches(7.05), Inches(1.0), Inches(0.35),
        f"{num}/{total}", font_size=9, color=MID_GRAY,
        alignment=PP_ALIGN.RIGHT
    )


def add_headline(slide, text, left=Inches(0.8), top=Inches(0.5),
                 width=Inches(11.7), font_size=30):
    """Standard slide headline."""
    add_textbox(slide, left, top, width, Inches(0.7), text,
                font_size=font_size, bold=True, color=CHARCOAL)


def add_footer(slide, text, top=Inches(6.8)):
    """Small footer text centered at the bottom."""
    add_textbox(
        slide, Inches(0.8), top, Inches(11.7), Inches(0.4),
        text, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.LEFT
    )


def add_green_dot_bullet(slide, left, top, width, text, font_size=14):
    """A bullet line with a green dot prefix."""
    dot_size = Inches(0.12)
    dot_top = top + Inches(0.06)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, dot_top, dot_size, dot_size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    add_textbox(slide, left + Inches(0.25), top, width - Inches(0.25),
                Inches(0.35), text, font_size=font_size, color=CHARCOAL)


# ── Slide Builders ───────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Thin green accent line near top
    add_accent_line(slide, Inches(5.5), Inches(2.2), Inches(2.3))

    # Tastyr IQ
    add_textbox(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(1.0),
                "Tastyr IQ", font_size=52, bold=True, color=CHARCOAL,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Tagline
    add_textbox(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(0.6),
                "Rate the dish. Not the restaurant.", font_size=22,
                bold=False, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Details
    add_textbox(slide, Inches(0), Inches(4.4), SLIDE_W, Inches(0.5),
                "AI-Powered Dish Intelligence  |  Seed Round  |  March 2026  |  Confidential",
                font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom accent line
    add_accent_line(slide, Inches(5.5), Inches(5.2), Inches(2.3))

    add_slide_number(slide, 1)


def slide_02_problem(prs):
    """Slide 2: The problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "A $371B industry built on the wrong data model")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    # 3 stat blocks side by side
    block_w = Inches(3.4)
    block_h = Inches(3.0)
    y = Inches(2.0)
    gap = Inches(0.55)
    x_start = Inches(1.2)

    stats = [
        ("$371B", "US restaurant industry"),
        ("0", "Platforms that score\nindividual dishes"),
        ("78%", "Consumers who check\nreviews before dining"),
    ]
    for i, (num, label) in enumerate(stats):
        x = x_start + i * (block_w + gap)
        add_metric_block(slide, x, y, block_w, block_h, num, label,
                         number_size=48, label_size=14)

    add_footer(slide, "Source: National Restaurant Association 2025")
    add_slide_number(slide, 2)


def slide_03_what_if(prs):
    """Slide 3: What if"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2),
                "What if you could score every dish\nin every restaurant \u2014 with AI?",
                font_size=36, bold=True, color=CHARCOAL,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=Pt(48))

    add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.6),
                "That's what we're building.", font_size=20,
                color=GREEN, alignment=PP_ALIGN.CENTER, bold=True)

    add_accent_line(slide, Inches(6.0), Inches(5.0), Inches(1.3))
    add_slide_number(slide, 3)


def slide_04_solution(prs):
    """Slide 4: Our solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "Tastyr IQ scores individual dishes 0\u2013100")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    # Left column: feature bullets
    bullets = [
        "Patent-pending AI scoring algorithm",
        "9 live food verticals",
        "Consumer + restaurant platform",
        "Data licensing infrastructure",
    ]
    y_start = Inches(2.0)
    for i, text in enumerate(bullets):
        add_green_dot_bullet(slide, Inches(1.0), y_start + Inches(i * 0.65),
                             Inches(5.0), text, font_size=16)

    # Right column: visual description block
    right_x = Inches(7.2)
    right_w = Inches(5.0)
    right_h = Inches(3.2)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Inches(1.8), right_w, right_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = OFF_WHITE
    shape.line.color.rgb = LIGHT_GRAY
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.04

    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)
    run = p.add_run()
    run.text = "TastyrScore"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = GREEN
    run.font.name = FONT

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(12)
    run2 = p2.add_run()
    run2.text = "87 / 100"
    run2.font.size = Pt(40)
    run2.font.bold = True
    run2.font.color.rgb = CHARCOAL
    run2.font.name = FONT

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "AI-scored dish cards with\nTastyrScore badges"
    run3.font.size = Pt(13)
    run3.font.color.rgb = MID_GRAY
    run3.font.name = FONT

    add_slide_number(slide, 4)


def slide_05_how_it_works(prs):
    """Slide 5: How it works"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "From reviews to scores in seconds")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    steps = [
        ("Aggregate", "Pull reviews, social\nmentions, critic ratings"),
        ("Analyze", "AI parses flavor, value,\npresentation, consistency"),
        ("Score", "Generate 0\u2013100\nTastyrScore per dish"),
        ("Deliver", "Consumer app +\nrestaurant dashboard"),
    ]

    step_w = Inches(2.4)
    gap = Inches(0.55)
    x_start = Inches(1.0)
    y_circle = Inches(2.4)
    y_title = Inches(3.2)
    y_desc = Inches(3.8)
    circle_size = Inches(0.55)

    for i, (title, desc) in enumerate(steps):
        x = x_start + i * (step_w + gap)
        cx = x + step_w / 2 - circle_size / 2

        # Circle
        add_step_circle(slide, cx, y_circle, circle_size, i + 1)

        # Arrow between steps
        if i < len(steps) - 1:
            arrow_x = x + step_w + Inches(0.05)
            add_arrow_right(slide, arrow_x, y_circle + Inches(0.12),
                            Inches(0.45), Inches(0.3))

        # Title
        add_textbox(slide, x, y_title, step_w, Inches(0.4),
                    title, font_size=18, bold=True, color=CHARCOAL,
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, x, y_desc, step_w, Inches(1.0),
                    desc, font_size=12, color=MID_GRAY,
                    alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 5)


def slide_06_market(prs):
    """Slide 6: Market opportunity"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "The market is massive and untapped")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    metrics = [
        ("$371B", "TAM", "US restaurant industry"),
        ("$23B", "SAM", "Restaurant tech & analytics SaaS"),
        ("$2.3B", "SOM", "Dish-level intelligence, Year 5"),
    ]

    block_w = Inches(3.4)
    block_h = Inches(3.0)
    gap = Inches(0.55)
    x_start = Inches(1.2)
    y = Inches(2.0)

    for i, (num, tag, label) in enumerate(metrics):
        x = x_start + i * (block_w + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, block_w, block_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = OFF_WHITE
        shape.line.color.rgb = LIGHT_GRAY
        shape.line.width = Pt(0.5)
        shape.adjustments[0] = 0.05

        tf = shape.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

        # Tag (TAM, SAM, SOM)
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(4)
        r0 = p0.add_run()
        r0.text = tag
        r0.font.size = Pt(14)
        r0.font.bold = True
        r0.font.color.rgb = GREEN
        r0.font.name = FONT

        # Number
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(6)
        r1 = p1.add_run()
        r1.text = num
        r1.font.size = Pt(44)
        r1.font.bold = True
        r1.font.color.rgb = CHARCOAL
        r1.font.name = FONT

        # Label
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(12)
        r2.font.color.rgb = MID_GRAY
        r2.font.name = FONT

    add_footer(slide, "Restaurant tech market growing at 12.5% CAGR through 2030")
    add_slide_number(slide, 6)


def slide_07_competitive(prs):
    """Slide 7: Competitive landscape"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "Nobody does what we do")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    # Table data
    headers = ["Feature", "Tastyr IQ", "Yelp", "Google Maps", "OpenTable"]
    rows = [
        ["Dish-level scores",    "\u2713", "\u2717", "\u2717", "\u2717"],
        ["AI scoring algorithm", "\u2713", "\u2717", "\u2717", "\u2717"],
        ["9 cuisine verticals",  "\u2713", "\u2717", "\u2717", "\u2717"],
        ["Restaurant analytics", "\u2713", "Paid",   "Limited", "\u2717"],
        ["Consumer discovery",   "\u2713", "\u2713", "\u2713",  "Partial"],
        ["Menu optimization",    "\u2713", "\u2717", "\u2717", "\u2717"],
        ["Data licensing",       "\u2713", "\u2717", "\u2717", "\u2717"],
    ]

    table_left = Inches(0.8)
    table_top = Inches(1.7)
    table_w = Inches(11.7)
    col_count = len(headers)
    row_count = len(rows) + 1

    table_shape = slide.shapes.add_table(row_count, col_count,
                                          table_left, table_top,
                                          table_w, Inches(4.5))
    table = table_shape.table

    # Column widths
    col_widths = [Inches(3.0), Inches(2.0), Inches(2.0), Inches(2.5), Inches(2.2)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Style header row
    for j, header_text in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = CHARCOAL

        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = header_text
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT

        # Remove borders
        _set_cell_border(cell, "thin")

    # Data rows
    for i, row_data in enumerate(rows):
        bg = WHITE if i % 2 == 0 else OFF_WHITE
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.size = Pt(12)
            run.font.name = FONT

            # Color checkmarks and x-marks
            if val == "\u2713":
                run.font.color.rgb = GREEN
                run.font.bold = True
                run.font.size = Pt(16)
            elif val == "\u2717":
                run.font.color.rgb = SOFT_RED
                run.font.size = Pt(16)
            else:
                run.font.color.rgb = CHARCOAL

            _set_cell_border(cell, "thin")

    add_slide_number(slide, 7)


def _set_cell_border(cell, style="thin"):
    """Set thin light gray borders on a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = tcPr.find(qn(edge))
        if ln is None:
            ln = tcPr.makeelement(qn(edge), {})
            tcPr.append(ln)
        ln.set("w", "6350")  # 0.5 pt
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        solidFill = ln.find(qn("a:solidFill"))
        if solidFill is None:
            solidFill = ln.makeelement(qn("a:solidFill"), {})
            ln.insert(0, solidFill)
        else:
            solidFill.clear()
        srgb = solidFill.makeelement(qn("a:srgbClr"), {"val": "E5E7EB"})
        solidFill.append(srgb)


def slide_08_business_model(prs):
    """Slide 8: Business model"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "Four revenue streams, one data platform")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    blocks = [
        ("40%", "B2C Subscriptions", "$4.99/mo"),
        ("35%", "B2B Restaurant SaaS", "$500/mo"),
        ("15%", "Data Licensing", "Enterprise"),
        ("10%", "Promoted Dishes", "CPC Model"),
    ]

    block_w = Inches(5.2)
    block_h = Inches(1.8)
    gap_x = Inches(0.6)
    gap_y = Inches(0.5)
    x_start = Inches(1.2)
    y_start = Inches(1.8)

    for i, (pct, title, price) in enumerate(blocks):
        col = i % 2
        row = i // 2
        x = x_start + col * (block_w + gap_x)
        y = y_start + row * (block_h + gap_y)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, block_w, block_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = OFF_WHITE
        shape.line.color.rgb = LIGHT_GRAY
        shape.line.width = Pt(0.5)
        shape.adjustments[0] = 0.05

        tf = shape.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

        # Percentage
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(2)
        r0 = p0.add_run()
        r0.text = pct
        r0.font.size = Pt(36)
        r0.font.bold = True
        r0.font.color.rgb = GREEN
        r0.font.name = FONT

        # Title
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(2)
        r1 = p1.add_run()
        r1.text = title
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = CHARCOAL
        r1.font.name = FONT

        # Price
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = price
        r2.font.size = Pt(13)
        r2.font.color.rgb = MID_GRAY
        r2.font.name = FONT

    add_slide_number(slide, 8)


def slide_09_traction(prs):
    """Slide 9: Traction & milestones"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "Built and validated in 8 months")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    metrics = [
        ("2,490+", "Dishes Indexed"),
        ("9", "Live Verticals"),
        ("93K", "Lines of Code"),
        ("1", "Patent Pending"),
        ("Live", "App Store"),
    ]

    block_w = Inches(2.0)
    block_h = Inches(2.8)
    gap = Inches(0.35)
    total_w = len(metrics) * block_w + (len(metrics) - 1) * gap
    x_start = (Inches(13.333) - total_w) / 2
    y = Inches(2.2)

    for i, (num, label) in enumerate(metrics):
        x = x_start + i * (block_w + gap)
        nc = GREEN if num != "Live" else AMBER
        add_metric_block(slide, x, y, block_w, block_h, num, label,
                         number_color=nc, number_size=38, label_size=12)

    add_slide_number(slide, 9)


def slide_10_gtm(prs):
    """Slide 10: Go-to-market strategy"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "Philadelphia first, then the nation")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    phases = [
        ("Phase 1: Launch", "Q2\u2013Q4 2026", "Philly launch, 1K users,\n10 restaurant partners"),
        ("Phase 2: Expand", "2027", "Tri-state expansion,\n8K users, data licensing"),
        ("Phase 3: Scale", "2028\u201329", "National rollout,\n100K users, enterprise deals"),
    ]

    block_w = Inches(3.4)
    block_h = Inches(3.2)
    gap = Inches(0.55)
    x_start = Inches(1.2)
    y = Inches(1.9)

    for i, (title, period, desc) in enumerate(phases):
        x = x_start + i * (block_w + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, block_w, block_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = OFF_WHITE
        shape.line.color.rgb = LIGHT_GRAY
        shape.line.width = Pt(0.5)
        shape.adjustments[0] = 0.05

        tf = shape.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(6)
        r0 = p0.add_run()
        r0.text = title
        r0.font.size = Pt(18)
        r0.font.bold = True
        r0.font.color.rgb = CHARCOAL
        r0.font.name = FONT

        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(10)
        r1 = p1.add_run()
        r1.text = period
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = GREEN
        r1.font.name = FONT

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = MID_GRAY
        r2.font.name = FONT

        # Arrow between blocks
        if i < len(phases) - 1:
            arrow_x = x + block_w + Inches(0.05)
            add_arrow_right(slide, arrow_x,
                            y + block_h / 2 - Inches(0.15),
                            Inches(0.45), Inches(0.3))

    add_slide_number(slide, 10)


def slide_11_social_chain_overview(prs):
    """Slide 11: Social chain strategy overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.4),
                "We don't cold-call investors.\nWe create organic market pull.",
                font_size=34, bold=True, color=CHARCOAL,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=Pt(46))

    add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.6),
                "8 interconnected influence chains that convert at 40%, not 2%",
                font_size=18, color=GREEN, alignment=PP_ALIGN.CENTER, bold=True)

    add_accent_line(slide, Inches(6.0), Inches(5.2), Inches(1.3))
    add_slide_number(slide, 11)


def _build_chain_slide(prs, slide_num, headline, steps):
    """Reusable chain strategy slide with vertical flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, headline)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    circle_size = Inches(0.45)
    y_start = Inches(1.9)
    step_h = Inches(0.85)

    for i, step_text in enumerate(steps):
        y = y_start + i * step_h

        # Circle
        add_step_circle(slide, Inches(1.2), y, circle_size, i + 1)

        # Down arrow between steps
        if i < len(steps) - 1:
            arrow_shape = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(1.32), y + circle_size + Inches(0.04),
                Inches(0.22), Inches(0.28)
            )
            arrow_shape.fill.solid()
            arrow_shape.fill.fore_color.rgb = LIGHT_GRAY
            arrow_shape.line.fill.background()

        # Text
        add_textbox(slide, Inches(1.9), y + Inches(0.04),
                    Inches(9.5), Inches(0.45),
                    step_text, font_size=15, color=CHARCOAL)

    add_slide_number(slide, slide_num)
    return slide


def slide_12_chain_a(prs):
    """Slide 12: Restaurateur Gateway"""
    _build_chain_slide(prs, 12, "Chain A: The Restaurateur Gateway", [
        "Food blogger posts dish scores \u2192 organic reach to restaurant community",
        "Restaurant owner sees coverage, gets curious about their own scores",
        "Restaurateur becomes early adopter, claims their menu",
        "At industry dinner, tells investor friends about the platform",
        "Investor hears from a trusted source, takes the call",
    ])


def slide_13_chain_b(prs):
    """Slide 13: FOMO Cascade"""
    _build_chain_slide(prs, 13, "Chain B: The FOMO Cascade", [
        "Seed 5\u201310 food influencers with exclusive beta access",
        "Coordinated posts drop in a 48-hour window",
        "Local media picks up the trend as a story",
        "Restaurateurs scramble to get rated before competitors",
        "Investors see market pull, not push \u2014 inbound interest spikes",
    ])


def slide_14_chain_c(prs):
    """Slide 14: University Pipeline"""
    _build_chain_slide(prs, 14, "Chain C: The University Pipeline", [
        "Faculty connection at Penn / Temple / Drexel",
        "Student beta testers sign up (200+ target)",
        "Usage data spike + campus buzz create social proof",
        "Faculty intro to university investor networks",
        "Alumni angel network presentation \u2192 warm capital",
    ])


def slide_15_chain_d(prs):
    """Slide 15: NJ-to-NYC Bridge"""
    _build_chain_slide(prs, 15, "Chain D: The NJ-to-NYC Bridge", [
        "Princeton Entrepreneurship Council demo",
        "NJ angel investors show interest, commit first checks",
        "NJ\u2013NYC co-investment bridge opens up",
        "NYC VCs see tri-state traction, join the round",
    ])


def slide_16_fomo_playbook(prs):
    """Slide 16: FOMO playbook"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "Creating urgency through scarcity and social proof")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    # Left column
    col_w = Inches(5.2)
    col_h = Inches(4.0)
    left_x = Inches(1.0)
    right_x = Inches(7.0)
    y = Inches(1.9)

    # Left block
    shape_l = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_x, y, col_w, col_h
    )
    shape_l.fill.solid()
    shape_l.fill.fore_color.rgb = OFF_WHITE
    shape_l.line.color.rgb = LIGHT_GRAY
    shape_l.line.width = Pt(0.5)
    shape_l.adjustments[0] = 0.04

    tf_l = shape_l.text_frame
    tf_l.word_wrap = True

    p0 = tf_l.paragraphs[0]
    p0.space_after = Pt(12)
    r0 = p0.add_run()
    r0.text = "Scarcity tactics"
    r0.font.size = Pt(18)
    r0.font.bold = True
    r0.font.color.rgb = CHARCOAL
    r0.font.name = FONT

    scarcity = [
        'Limited beta slots (50 restaurants)',
        '"Founding Partner" tier closes permanently',
        'Exclusive data access window for early investors',
        'Capped round with first-come allocation',
    ]
    for b in scarcity:
        p = tf_l.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"\u2022  {b}"
        run.font.size = Pt(13)
        run.font.color.rgb = CHARCOAL
        run.font.name = FONT

    # Right block
    shape_r = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, y, col_w, col_h
    )
    shape_r.fill.solid()
    shape_r.fill.fore_color.rgb = OFF_WHITE
    shape_r.line.color.rgb = LIGHT_GRAY
    shape_r.line.width = Pt(0.5)
    shape_r.adjustments[0] = 0.04

    tf_r = shape_r.text_frame
    tf_r.word_wrap = True

    p0r = tf_r.paragraphs[0]
    p0r.space_after = Pt(12)
    r0r = p0r.add_run()
    r0r.text = "Social proof manufacturing"
    r0r.font.size = Pt(18)
    r0r.font.bold = True
    r0r.font.color.rgb = CHARCOAL
    r0r.font.name = FONT

    proof = [
        'Coordinated influencer content drops',
        'Restaurant partner press announcements',
        'University engagement metrics published',
        'Investor interest ticker on landing page',
    ]
    for b in proof:
        p = tf_r.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"\u2022  {b}"
        run.font.size = Pt(13)
        run.font.color.rgb = CHARCOAL
        run.font.name = FONT

    add_slide_number(slide, 16)


def slide_17_team(prs):
    """Slide 17: The team"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "The team")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    profiles = [
        {
            "name": "Nik Amin",
            "title": "CEO & Co-Founder",
            "details": [
                "Bay Area \u2014 Product & Engineering",
                "Built entire platform (93K lines)",
                "Patent-pending AI scoring algorithm",
                "Full-stack: React, Python, ML pipeline",
            ],
        },
        {
            "name": "Ali Mabsoute",
            "title": "CMO & Co-Founder",
            "details": [
                "UPenn grad, 15 years in Philadelphia",
                "Hyatt, Citi, American Express",
                "Deep Philadelphia market knowledge",
                "Growth strategy & investor relations",
            ],
        },
    ]

    block_w = Inches(5.2)
    block_h = Inches(4.0)
    gap = Inches(0.6)
    x_start = Inches(1.2)
    y = Inches(1.9)

    for i, profile in enumerate(profiles):
        x = x_start + i * (block_w + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, block_w, block_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = OFF_WHITE
        shape.line.color.rgb = LIGHT_GRAY
        shape.line.width = Pt(0.5)
        shape.adjustments[0] = 0.04

        tf = shape.text_frame
        tf.word_wrap = True

        # Name
        p0 = tf.paragraphs[0]
        p0.space_after = Pt(4)
        r0 = p0.add_run()
        r0.text = profile["name"]
        r0.font.size = Pt(24)
        r0.font.bold = True
        r0.font.color.rgb = CHARCOAL
        r0.font.name = FONT

        # Title
        p1 = tf.add_paragraph()
        p1.space_after = Pt(14)
        r1 = p1.add_run()
        r1.text = profile["title"]
        r1.font.size = Pt(15)
        r1.font.bold = True
        r1.font.color.rgb = GREEN
        r1.font.name = FONT

        # Detail bullets
        for detail in profile["details"]:
            p = tf.add_paragraph()
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = f"\u2022  {detail}"
            run.font.size = Pt(13)
            run.font.color.rgb = CHARCOAL
            run.font.name = FONT

    add_slide_number(slide, 17)


def slide_18_financials(prs):
    """Slide 18: Financial projections"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "Path to $18M ARR by 2029")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    # Top row: ARR metrics
    top_metrics = [
        ("$120K", "2026 ARR"),
        ("$1.2M", "2027 ARR"),
        ("$18M", "2029 ARR"),
    ]
    # Bottom row: efficiency metrics
    bot_metrics = [
        ("70%+", "Gross Margin"),
        ("3 mo", "CAC Payback"),
        ("6.5x", "LTV:CAC"),
    ]

    block_w = Inches(3.4)
    block_h = Inches(1.9)
    gap = Inches(0.55)
    x_start = Inches(1.2)

    for row_i, metrics in enumerate([top_metrics, bot_metrics]):
        y = Inches(1.8) + row_i * (block_h + Inches(0.4))
        for j, (num, label) in enumerate(metrics):
            x = x_start + j * (block_w + gap)
            nc = GREEN if row_i == 0 else CHARCOAL
            add_metric_block(slide, x, y, block_w, block_h, num, label,
                             number_color=nc, number_size=36, label_size=12)

    add_footer(slide, "Path to profitability in Q4 2028")
    add_slide_number(slide, 18)


def slide_19_the_ask(prs):
    """Slide 19: The ask"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "Raising $500K\u2013$1.5M seed")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    # Valuation callout
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5),
                "$3M\u2013$6M pre-money valuation", font_size=18,
                bold=True, color=GREEN, alignment=PP_ALIGN.LEFT)

    # 4 use-of-funds blocks in a row
    blocks = [
        ("40%", "Engineering & AI"),
        ("25%", "Growth & Marketing"),
        ("20%", "Business Development"),
        ("15%", "Operations"),
    ]

    block_w = Inches(2.6)
    block_h = Inches(3.0)
    gap = Inches(0.4)
    total_w = len(blocks) * block_w + (len(blocks) - 1) * gap
    x_start = (Inches(13.333) - total_w) / 2
    y = Inches(2.5)

    for i, (pct, label) in enumerate(blocks):
        x = x_start + i * (block_w + gap)
        add_metric_block(slide, x, y, block_w, block_h, pct, label,
                         number_color=GREEN, number_size=40, label_size=13)

    add_slide_number(slide, 19)


def slide_20_action_plan(prs):
    """Slide 20: Week-by-week action plan"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_headline(slide, "12-week sprint to close")
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5))

    phases = [
        ("Week 1\u20132", "Foundation", "Recruit influencers, onboard restaurants,\nsubmit accelerator applications"),
        ("Week 3\u20134", "Momentum", "Coordinated content drops, media outreach,\nuniversity beta launch"),
        ("Week 5\u20138", "Acceleration", "10+ investor meetings, shift from\npush to pull, FOMO builds"),
        ("Week 9\u201312", "Close", "Term sheets, strategic investors,\nfunding announcement"),
    ]

    block_w = Inches(2.6)
    block_h = Inches(3.5)
    gap = Inches(0.4)
    total_w = len(phases) * block_w + (len(phases) - 1) * gap
    x_start = (Inches(13.333) - total_w) / 2
    y = Inches(1.8)

    for i, (weeks, title, desc) in enumerate(phases):
        x = x_start + i * (block_w + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, block_w, block_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = OFF_WHITE
        shape.line.color.rgb = LIGHT_GRAY
        shape.line.width = Pt(0.5)
        shape.adjustments[0] = 0.05

        tf = shape.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.space_after = Pt(6)
        r0 = p0.add_run()
        r0.text = weeks
        r0.font.size = Pt(14)
        r0.font.bold = True
        r0.font.color.rgb = GREEN
        r0.font.name = FONT

        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(10)
        r1 = p1.add_run()
        r1.text = title
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = CHARCOAL
        r1.font.name = FONT

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = MID_GRAY
        r2.font.name = FONT

        # Arrow between
        if i < len(phases) - 1:
            arrow_x = x + block_w + Inches(0.02)
            add_arrow_right(slide, arrow_x,
                            y + block_h / 2 - Inches(0.15),
                            Inches(0.35), Inches(0.3))

    add_slide_number(slide, 20)


def slide_21_contact(prs):
    """Slide 21: Contact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_accent_line(slide, Inches(5.5), Inches(2.0), Inches(2.3))

    add_textbox(slide, Inches(0), Inches(2.3), SLIDE_W, Inches(0.9),
                "Tastyr IQ", font_size=48, bold=True, color=CHARCOAL,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, Inches(0), Inches(3.3), SLIDE_W, Inches(0.6),
                "Rate the dish. Not the restaurant.", font_size=20,
                color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0), Inches(4.2), SLIDE_W, Inches(0.5),
                "nik@tastyiq.com  |  tastyiq.com", font_size=16,
                color=GREEN, alignment=PP_ALIGN.CENTER, bold=True)

    add_textbox(slide, Inches(0), Inches(5.0), SLIDE_W, Inches(0.4),
                "Patent Pending  |  App Store Live  |  Philadelphia, PA",
                font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(5.6), Inches(2.3))
    add_slide_number(slide, 21)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()

    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build all 21 slides
    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_what_if(prs)
    slide_04_solution(prs)
    slide_05_how_it_works(prs)
    slide_06_market(prs)
    slide_07_competitive(prs)
    slide_08_business_model(prs)
    slide_09_traction(prs)
    slide_10_gtm(prs)
    slide_11_social_chain_overview(prs)
    slide_12_chain_a(prs)
    slide_13_chain_b(prs)
    slide_14_chain_c(prs)
    slide_15_chain_d(prs)
    slide_16_fomo_playbook(prs)
    slide_17_team(prs)
    slide_18_financials(prs)
    slide_19_the_ask(prs)
    slide_20_action_plan(prs)
    slide_21_contact(prs)

    # Save
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "tastyr-iq-outreach-deck.pptx")
    prs.save(out_path)
    print(f"Deck saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
