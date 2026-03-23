"""Build 4 Tastyr IQ pitch deck variations from the Yelp template PPTX."""
import os
from pptx import Presentation

TEMPLATE = r"C:\Users\alima\Downloads\yelp_app_pitch_deck.pptx"
OUT_DIR = r"C:\Users\alima\tastyr-iq\strategy"


def replace_shape_text(shape, new_text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_text
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ""
    for pi in range(1, len(tf.paragraphs)):
        for r in tf.paragraphs[pi].runs:
            r.text = ""


def set_multiline(shape, lines):
    if not shape.has_text_frame:
        return
    paras = shape.text_frame.paragraphs
    for i, line in enumerate(lines):
        if i < len(paras) and paras[i].runs:
            paras[i].runs[0].text = line
            for r in paras[i].runs[1:]:
                r.text = ""
    for i in range(len(lines), len(paras)):
        for r in paras[i].runs:
            r.text = ""


def get_shapes_by_name(slide):
    return {s.name: s for s in slide.shapes}


def build_deck(filename, v):
    prs = Presentation(TEMPLATE)
    slides = list(prs.slides)

    # SLIDE 1: Title/Cover
    sh = get_shapes_by_name(slides[0])
    replace_shape_text(sh['Text 2'], v['s1_brand'])
    replace_shape_text(sh['Text 3'], v['s1_tagline'])
    replace_shape_text(sh['Text 5'], v['s1_subtitle'])
    replace_shape_text(sh['Text 6'], v['s1_meta'])

    # SLIDE 2: The Problem
    sh = get_shapes_by_name(slides[1])
    replace_shape_text(sh['Text 1'], v['s2_title'])
    replace_shape_text(sh['Text 2'], v['s2_headline'])
    replace_shape_text(sh['Text 4'], v['s2_card1_title'])
    replace_shape_text(sh['Text 5'], v['s2_card1_body'])
    replace_shape_text(sh['Text 7'], v['s2_card2_title'])
    replace_shape_text(sh['Text 8'], v['s2_card2_body'])
    replace_shape_text(sh['Text 10'], v['s2_card3_title'])
    replace_shape_text(sh['Text 11'], v['s2_card3_body'])
    replace_shape_text(sh['Text 13'], v['s2_stat'])
    set_multiline(sh['Text 14'], v['s2_stat_desc'])
    replace_shape_text(sh['Text 15'], v['s2_source'])

    # SLIDE 3: Our Solution
    sh = get_shapes_by_name(slides[2])
    replace_shape_text(sh['Text 1'], 'Our Solution')
    replace_shape_text(sh['Text 2'], v['s3_desc'])
    replace_shape_text(sh['Text 4'], v['s3_f1_title'])
    replace_shape_text(sh['Text 5'], v['s3_f1_body'])
    replace_shape_text(sh['Text 7'], v['s3_f2_title'])
    replace_shape_text(sh['Text 8'], v['s3_f2_body'])
    replace_shape_text(sh['Text 10'], v['s3_f3_title'])
    replace_shape_text(sh['Text 11'], v['s3_f3_body'])
    replace_shape_text(sh['Text 13'], v['s3_f4_title'])
    replace_shape_text(sh['Text 14'], v['s3_f4_body'])

    # SLIDE 4: Product Overview
    sh = get_shapes_by_name(slides[3])
    replace_shape_text(sh['Text 1'], 'Product Overview')
    replace_shape_text(sh['Text 3'], v['s4_brand'])
    replace_shape_text(sh['Text 5'], v['s4_search'])
    replace_shape_text(sh['Text 7'], v['s4_item_title'])
    # Dish list items (left side)
    replace_shape_text(sh['Text 8'], v.get('s4_dish1_sub', 'Thai  \u2605 96/100'))
    replace_shape_text(sh['Text 10'], v.get('s4_dish1_badge', '96'))
    replace_shape_text(sh['Text 12'], v.get('s4_dish2_name', 'Truffle Mushroom Risotto'))
    replace_shape_text(sh['Text 13'], v.get('s4_dish2_sub', 'Italian  \u2605 94/100'))
    replace_shape_text(sh['Text 15'], v.get('s4_dish2_badge', '94'))
    replace_shape_text(sh['Text 17'], v.get('s4_dish3_name', 'Nashville Hot Chicken'))
    replace_shape_text(sh['Text 18'], v.get('s4_dish3_sub', 'Southern  \u2605 91/100'))
    replace_shape_text(sh['Text 20'], v.get('s4_dish3_badge', '91'))
    # Feature cards (right side)
    replace_shape_text(sh['Text 21'], v.get('s4_feat1_title', 'Patent-Pending Dish Scores'))
    replace_shape_text(sh['Text 22'], v.get('s4_feat1_body', 'Every dish gets a 0-100 TastyrScore powered by our proprietary AI algorithm across 9 food verticals.'))
    replace_shape_text(sh['Text 23'], v.get('s4_feat2_title', 'Cuisine Vertical AI'))
    replace_shape_text(sh['Text 24'], v.get('s4_feat2_body', 'Specialized models for each cuisine — what makes great sushi is different from great BBQ.'))
    replace_shape_text(sh['Text 25'], v.get('s4_feat3_title', 'Real-Time Dish Rankings'))
    replace_shape_text(sh['Text 26'], v.get('s4_feat3_body', 'Live trending dishes, new additions, and neighborhood-level dish intelligence updated continuously.'))
    replace_shape_text(sh['Text 27'], v.get('s4_feat4_title', 'Restaurant Dashboard'))
    replace_shape_text(sh['Text 28'], v.get('s4_feat4_body', 'Analytics for restaurants — dish performance, competitive benchmarking, menu optimization insights.'))

    # SLIDE 5: Market Opportunity
    sh = get_shapes_by_name(slides[4])
    replace_shape_text(sh['Text 1'], 'Market Opportunity')
    replace_shape_text(sh['Text 3'], v['s5_tam_label'])
    replace_shape_text(sh['Text 4'], v['s5_tam_val'])
    set_multiline(sh['Text 5'], v['s5_tam_desc'])
    replace_shape_text(sh['Text 7'], v['s5_sam_label'])
    replace_shape_text(sh['Text 8'], v['s5_sam_val'])
    set_multiline(sh['Text 9'], v['s5_sam_desc'])
    replace_shape_text(sh['Text 11'], v['s5_som_label'])
    replace_shape_text(sh['Text 12'], v['s5_som_val'])
    set_multiline(sh['Text 13'], v['s5_som_desc'])
    replace_shape_text(sh['Text 14'], v['s5_footer'])

    # SLIDE 6: Business Model
    sh = get_shapes_by_name(slides[5])
    replace_shape_text(sh['Text 1'], 'Business Model')
    replace_shape_text(sh['Text 4'], v['s6_pct1'])
    replace_shape_text(sh['Text 5'], v['s6_rev1'])
    set_multiline(sh['Text 6'], v['s6_rev1_desc'])
    replace_shape_text(sh['Text 9'], v['s6_pct2'])
    replace_shape_text(sh['Text 10'], v['s6_rev2'])
    set_multiline(sh['Text 11'], v['s6_rev2_desc'])
    replace_shape_text(sh['Text 14'], v['s6_pct3'])
    replace_shape_text(sh['Text 15'], v['s6_rev3'])
    set_multiline(sh['Text 16'], v['s6_rev3_desc'])
    replace_shape_text(sh['Text 19'], v['s6_pct4'])
    replace_shape_text(sh['Text 20'], v['s6_rev4'])
    set_multiline(sh['Text 21'], v['s6_rev4_desc'])

    # SLIDE 7: Competitive Landscape
    sh = get_shapes_by_name(slides[6])
    replace_shape_text(sh['Text 1'], 'Competitive Landscape')
    replace_shape_text(sh['Text 2'], v['s7_desc'])
    # Replace competitive table data
    for s in slides[6].shapes:
        if s.has_table:
            tbl = s.table
            comp_rows = v.get('s7_table', [
                ['Feature', 'Tastyr IQ', 'Yelp', 'Google Maps', 'OpenTable'],
                ['Dish-Level Scores', '\u2713', '\u2717', '\u2717', '\u2717'],
                ['AI Scoring Algorithm', '\u2713', '\u2717', '\u2717', '\u2717'],
                ['9 Cuisine Verticals', '\u2713', '\u2717', '\u2717', '\u2717'],
                ['Restaurant Dashboard', '\u2713', 'Paid Only', 'Limited', '\u2717'],
                ['Consumer Discovery', '\u2713', '\u2713', '\u2713', 'Partial'],
                ['Menu Optimization', '\u2713', '\u2717', '\u2717', '\u2717'],
                ['Data Licensing', '\u2713', '\u2717', '\u2717', '\u2717'],
            ])
            for ri, row_data in enumerate(comp_rows):
                if ri < len(list(tbl.rows)):
                    row = list(tbl.rows)[ri]
                    for ci, cell_text in enumerate(row_data):
                        if ci < len(tbl.columns):
                            cell = row.cells[ci]
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    run.text = ""
                            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                                cell.text_frame.paragraphs[0].runs[0].text = cell_text
                            elif cell.text_frame.paragraphs:
                                cell.text_frame.paragraphs[0].text = cell_text
            break

    # SLIDE 8: Go-to-Market Strategy
    sh = get_shapes_by_name(slides[7])
    replace_shape_text(sh['Text 1'], 'Go-to-Market Strategy')
    replace_shape_text(sh['Text 4'], v['s8_p1_phase'])
    replace_shape_text(sh['Text 5'], v['s8_p1_name'])
    replace_shape_text(sh['Text 6'], v['s8_p1_time'])
    set_multiline(sh['Text 7'], v['s8_p1_items'])
    replace_shape_text(sh['Text 10'], v['s8_p2_phase'])
    replace_shape_text(sh['Text 11'], v['s8_p2_name'])
    replace_shape_text(sh['Text 12'], v['s8_p2_time'])
    set_multiline(sh['Text 13'], v['s8_p2_items'])
    replace_shape_text(sh['Text 16'], v['s8_p3_phase'])
    replace_shape_text(sh['Text 17'], v['s8_p3_name'])
    replace_shape_text(sh['Text 18'], v['s8_p3_time'])
    set_multiline(sh['Text 19'], v['s8_p3_items'])

    # SLIDE 9: Traction & Milestones
    sh = get_shapes_by_name(slides[8])
    replace_shape_text(sh['Text 1'], 'Traction & Milestones')
    replace_shape_text(sh['Text 2'], v['s9_headline'])
    replace_shape_text(sh['Text 4'], v['s9_s1_val'])
    replace_shape_text(sh['Text 5'], v['s9_s1_label'])
    replace_shape_text(sh['Text 7'], v['s9_s2_val'])
    replace_shape_text(sh['Text 8'], v['s9_s2_label'])
    replace_shape_text(sh['Text 10'], v['s9_s3_val'])
    replace_shape_text(sh['Text 11'], v['s9_s3_label'])
    replace_shape_text(sh['Text 13'], v['s9_s4_val'])
    replace_shape_text(sh['Text 14'], v['s9_s4_label'])
    replace_shape_text(sh['Text 16'], v['s9_s5_val'])
    replace_shape_text(sh['Text 17'], v['s9_s5_label'])

    # SLIDE 10: The Team
    sh = get_shapes_by_name(slides[9])
    replace_shape_text(sh['Text 1'], 'The Team')
    replace_shape_text(sh['Text 4'], v['s10_m1_init'])
    replace_shape_text(sh['Text 5'], v['s10_m1_name'])
    replace_shape_text(sh['Text 6'], v['s10_m1_role'])
    set_multiline(sh['Text 7'], v['s10_m1_desc'])
    replace_shape_text(sh['Text 10'], v['s10_m2_init'])
    replace_shape_text(sh['Text 11'], v['s10_m2_name'])
    replace_shape_text(sh['Text 12'], v['s10_m2_role'])
    set_multiline(sh['Text 13'], v['s10_m2_desc'])
    replace_shape_text(sh['Text 16'], v['s10_m3_init'])
    replace_shape_text(sh['Text 17'], v['s10_m3_name'])
    replace_shape_text(sh['Text 18'], v['s10_m3_role'])
    set_multiline(sh['Text 19'], v['s10_m3_desc'])
    replace_shape_text(sh['Text 22'], v['s10_m4_init'])
    replace_shape_text(sh['Text 23'], v['s10_m4_name'])
    replace_shape_text(sh['Text 24'], v['s10_m4_role'])
    set_multiline(sh['Text 25'], v['s10_m4_desc'])
    replace_shape_text(sh['Text 26'], v['s10_advisory'])

    # SLIDE 11: Financial Projections
    sh = get_shapes_by_name(slides[10])
    replace_shape_text(sh['Text 1'], 'Financial Projections')
    replace_shape_text(sh['Text 3'], v['s11_v1'])
    replace_shape_text(sh['Text 4'], v['s11_l1'])
    replace_shape_text(sh['Text 6'], v['s11_v2'])
    replace_shape_text(sh['Text 7'], v['s11_l2'])
    replace_shape_text(sh['Text 9'], v['s11_v3'])
    replace_shape_text(sh['Text 10'], v['s11_l3'])
    replace_shape_text(sh['Text 12'], v['s11_v4'])
    replace_shape_text(sh['Text 13'], v['s11_l4'])
    replace_shape_text(sh['Text 15'], v['s11_v5'])
    replace_shape_text(sh['Text 16'], v['s11_l5'])
    replace_shape_text(sh['Text 18'], v['s11_v6'])
    replace_shape_text(sh['Text 19'], v['s11_l6'])
    replace_shape_text(sh['Text 20'], v['s11_footer'])

    # SLIDE 12: The Ask
    sh = get_shapes_by_name(slides[11])
    replace_shape_text(sh['Text 1'], 'The Ask')
    replace_shape_text(sh['Text 2'], v['s12_amount'])
    replace_shape_text(sh['Text 3'], v['s12_valuation'])
    replace_shape_text(sh['Text 4'], v['s12_uof'])
    replace_shape_text(sh['Text 5'], v['s12_pct1'])
    replace_shape_text(sh['Text 6'], v['s12_uof1'])
    replace_shape_text(sh['Text 7'], v['s12_uof1d'])
    replace_shape_text(sh['Text 8'], v['s12_pct2'])
    replace_shape_text(sh['Text 9'], v['s12_uof2'])
    replace_shape_text(sh['Text 10'], v['s12_uof2d'])
    replace_shape_text(sh['Text 11'], v['s12_pct3'])
    replace_shape_text(sh['Text 12'], v['s12_uof3'])
    replace_shape_text(sh['Text 13'], v['s12_uof3d'])
    replace_shape_text(sh['Text 14'], v['s12_pct4'])
    replace_shape_text(sh['Text 15'], v['s12_uof4'])
    replace_shape_text(sh['Text 16'], v['s12_uof4d'])
    replace_shape_text(sh['Text 18'], v['s12_cta'])

    out_path = os.path.join(OUT_DIR, filename)
    prs.save(out_path)
    print(f"SAVED: {out_path}")


# ============================================================
# VARIATION 1: INVESTOR FOCUS
# ============================================================
V1 = {
    "s1_brand": "Tastyr IQ",
    "s1_tagline": "Rate the dish. Not the restaurant.",
    "s1_subtitle": "AI-Powered Dish Intelligence Platform",
    "s1_meta": "Seed Round  |  March 2026  |  Confidential",
    "s2_title": "The Problem",
    "s2_headline": "A $371B industry with zero dish-level data.",
    "s2_card1_title": "Restaurants Are Not Dishes",
    "s2_card1_body": "A 4-star restaurant can serve 2-star dishes. Today's platforms rate the venue, not the food.",
    "s2_card2_title": "No Objective Dish Scores",
    "s2_card2_body": "Consumers scroll through 500-word reviews hoping to find one mention of what they want to order.",
    "s2_card3_title": "Restaurants Fly Blind on Menus",
    "s2_card3_body": "Operators have POS data but zero insight into which dishes actually drive loyalty and return visits.",
    "s2_stat": "$371B",
    "s2_stat_desc": ["US restaurant industry", "with no dish-level", "data infrastructure"],
    "s2_source": "Source: National Restaurant Association 2025",
    "s3_desc": "Tastyr IQ is the first AI platform that scores individual dishes 0-100 using patent-pending algorithms across 9 food verticals.",
    "s3_f1_title": "Patent-Pending Dish Scoring",
    "s3_f1_body": "Proprietary AI algorithm analyzes flavor, presentation, value, and consistency to generate a 0-100 TastyrScore for every dish.",
    "s3_f2_title": "9 Live Food Verticals",
    "s3_f2_body": "From sushi to BBQ to vegan, our vertical-specific models understand what makes each cuisine category excellent.",
    "s3_f3_title": "Consumer + Restaurant Platform",
    "s3_f3_body": "Consumers discover top-scored dishes. Restaurants get analytics dashboards showing dish-level performance.",
    "s3_f4_title": "Data Licensing Infrastructure",
    "s3_f4_body": "Dish intelligence data licensed to delivery platforms, CPG brands, and food industry analysts.",
    "s4_brand": "TASTYR IQ", "s4_search": "Search dishes near you...", "s4_item_title": "Spicy Tuna Crispy Rice",
    "s4_dish1_sub": "Japanese  \u2605 96/100", "s4_dish1_badge": "96",
    "s4_dish2_name": "Truffle Mushroom Risotto", "s4_dish2_sub": "Italian  \u2605 94/100", "s4_dish2_badge": "94",
    "s4_dish3_name": "Nashville Hot Chicken", "s4_dish3_sub": "Southern  \u2605 91/100", "s4_dish3_badge": "91",
    "s4_feat1_title": "Patent-Pending Dish Scores", "s4_feat1_body": "Every dish gets a 0-100 TastyrScore using AI that analyzes flavor, presentation, value, and consistency.",
    "s4_feat2_title": "9 Cuisine Vertical AI", "s4_feat2_body": "Specialized models for each food category — sushi, BBQ, vegan, Italian, and more.",
    "s4_feat3_title": "Real-Time Dish Rankings", "s4_feat3_body": "Live trending dishes, new additions, and neighborhood-level intelligence updated continuously.",
    "s4_feat4_title": "Restaurant Dashboard", "s4_feat4_body": "Free analytics for restaurants — dish performance, competitive benchmarking, menu optimization insights.",
    "s5_tam_label": "TAM", "s5_tam_val": "$371B", "s5_tam_desc": ["US restaurant industry", "total addressable market"],
    "s5_sam_label": "SAM", "s5_sam_val": "$23B", "s5_sam_desc": ["Restaurant tech & analytics", "SaaS market"],
    "s5_som_label": "SOM", "s5_som_val": "$2.3B", "s5_som_desc": ["Dish-level intelligence", "Year 5 target"],
    "s5_footer": "Restaurant tech market growing at 12.5% CAGR through 2030",
    "s6_pct1": "40%", "s6_rev1": "B2C Subscriptions", "s6_rev1_desc": ["Premium dish discovery, personalized recs", "$4.99/mo per user"],
    "s6_pct2": "35%", "s6_rev2": "B2B Restaurant SaaS", "s6_rev2_desc": ["Dish analytics dashboard, menu optimization", "$500/mo per restaurant"],
    "s6_pct3": "15%", "s6_rev3": "Data Licensing", "s6_rev3_desc": ["Dish intelligence for delivery platforms & CPG", "Enterprise contracts"],
    "s6_pct4": "10%", "s6_rev4": "Promoted Dishes", "s6_rev4_desc": ["Restaurants promote top-scored dishes", "CPC model"],
    "s7_desc": "Tastyr IQ is the only platform that scores individual dishes with AI. No competitor offers dish-level intelligence.",
    "s8_p1_phase": "Phase 1", "s8_p1_name": "Launch", "s8_p1_time": "Q2-Q4 2026",
    "s8_p1_items": ["Philadelphia launch city", "1,000 premium subscribers", "First 10 restaurant dashboards", "Founding user cohort campaign"],
    "s8_p2_phase": "Phase 2", "s8_p2_name": "Expand", "s8_p2_time": "2027",
    "s8_p2_items": ["Tri-state expansion (NYC, NJ)", "8,000 premium users", "80 restaurant dashboards", "First data licensing deals"],
    "s8_p3_phase": "Phase 3", "s8_p3_name": "Scale", "s8_p3_time": "2028-2029",
    "s8_p3_items": ["National expansion to top 25 cities", "100K premium users", "800 restaurant dashboards", "Enterprise data licensing"],
    "s9_headline": "Built and validated in 8 months",
    "s9_s1_val": "2,490+", "s9_s1_label": "Dishes Indexed & Scored",
    "s9_s2_val": "9", "s9_s2_label": "Live Food Verticals",
    "s9_s3_val": "93K", "s9_s3_label": "Lines of Production Code",
    "s9_s4_val": "1", "s9_s4_label": "Patent Pending (Scoring Algo)",
    "s9_s5_val": "Live", "s9_s5_label": "TestFlight Beta + App Store",
    "s10_m1_init": "NA", "s10_m1_name": "Nik Amin", "s10_m1_role": "CEO & Co-Founder",
    "s10_m1_desc": ["Bay Area, Product & Engineering", "Built entire platform, 93K lines of code"],
    "s10_m2_init": "AM", "s10_m2_name": "Ali Mabsoute", "s10_m2_role": "CMO & Co-Founder",
    "s10_m2_desc": ["UPenn Grad, 15 Years in Philadelphia", "Hyatt, Citi, American Express alumni"],
    "s10_m3_init": "", "s10_m3_name": "", "s10_m3_role": "", "s10_m3_desc": ["", ""],
    "s10_m4_init": "", "s10_m4_name": "", "s10_m4_role": "", "s10_m4_desc": ["", ""],
    "s10_advisory": "Deep domain expertise in food tech, restaurant operations, and consumer product growth.",
    "s11_v1": "$120K", "s11_l1": "2026 ARR",
    "s11_v2": "$1.2M", "s11_l2": "2027 ARR",
    "s11_v3": "$18M", "s11_l3": "2029 ARR (Target)",
    "s11_v4": "70%+", "s11_l4": "Gross Margin",
    "s11_v5": "3 months", "s11_l5": "CAC Payback (B2C)",
    "s11_v6": "6.5x", "s11_l6": "LTV:CAC Ratio",
    "s11_footer": "Path to profitability in Q4 2028  |  Conservative projections, multiple revenue streams",
    "s12_amount": "$500K - $1.5M Seed",
    "s12_valuation": "at a $3M-$6M pre-money valuation",
    "s12_uof": "Use of Funds",
    "s12_pct1": "40%", "s12_uof1": "Engineering & AI", "s12_uof1d": "Scoring models, infrastructure, vertical expansion",
    "s12_pct2": "25%", "s12_uof2": "Growth & Marketing", "s12_uof2d": "User acquisition, influencer campaigns",
    "s12_pct3": "20%", "s12_uof3": "Business Development", "s12_uof3d": "Restaurant partnerships, data licensing",
    "s12_pct4": "15%", "s12_uof4": "Operations", "s12_uof4d": "Legal, IP protection, admin",
    "s12_cta": "nik@tastyiq.com  |  tastyiq.com  |  The dish-level data layer the food industry has been missing.",
}

# ============================================================
# VARIATION 2: FOMO / MOMENTUM
# ============================================================
V2 = {
    "s1_brand": "Tastyr IQ",
    "s1_tagline": "The next category-defining food tech platform.",
    "s1_subtitle": "Dish Intelligence Is Coming. We're Building It First.",
    "s1_meta": "Seed Round  |  March 2026  |  Confidential",
    "s2_title": "The Problem",
    "s2_headline": "Someone will own dish-level intelligence. The question is who.",
    "s2_card1_title": "A Massive Blind Spot",
    "s2_card1_body": "The $371B restaurant industry has never had a data layer for individual dishes. That gap is closing fast.",
    "s2_card2_title": "Big Tech Is Circling",
    "s2_card2_body": "Google, Yelp, and DoorDash have the resources but haven't moved yet. The first-mover window is NOW.",
    "s2_card3_title": "Consumer Demand Is Exploding",
    "s2_card3_body": "TikTok food content gets 40B+ views. People want to know WHAT to eat, not just WHERE to eat.",
    "s2_stat": "0",
    "s2_stat_desc": ["platforms today score", "individual dishes", "with AI"],
    "s2_source": "The category is wide open.",
    "s3_desc": "Tastyr IQ is building the Rotten Tomatoes for Food. AI-powered scores for every dish, in every restaurant, in every city.",
    "s3_f1_title": "First-Mover Advantage",
    "s3_f1_body": "Patent-pending scoring algorithm. No competitor has filed. We own the intellectual property for dish intelligence.",
    "s3_f2_title": "App Store Live TODAY",
    "s3_f2_body": "Not a concept. Not a prototype. Tastyr IQ is approved and live on the App Store. Users can download right now.",
    "s3_f3_title": "Network Effects Compound",
    "s3_f3_body": "Every dish scored, every review analyzed, every restaurant added makes the platform exponentially more valuable.",
    "s3_f4_title": "Winner-Take-All Dynamics",
    "s3_f4_body": "Like Rotten Tomatoes for movies, the first credible dish scoring platform becomes the default. Second place is irrelevant.",
    "s4_brand": "TASTYR IQ", "s4_search": "Find the best dishes near you...", "s4_item_title": "Birria Tacos",
    "s4_dish1_sub": "Mexican  \u2605 97/100", "s4_dish1_badge": "97",
    "s4_dish2_name": "Khao Soi Curry", "s4_dish2_sub": "Thai  \u2605 95/100", "s4_dish2_badge": "95",
    "s4_dish3_name": "Wood-Fired Margherita", "s4_dish3_sub": "Italian  \u2605 93/100", "s4_dish3_badge": "93",
    "s4_feat1_title": "First-Mover Dish Scoring", "s4_feat1_body": "Patent-pending AI scores every dish 0-100. No competitor has filed. We own the IP.",
    "s4_feat2_title": "App Store Live NOW", "s4_feat2_body": "Not a concept. Users can download today. 2,490+ dishes already scored and live.",
    "s4_feat3_title": "Network Effects Lock-In", "s4_feat3_body": "Every dish scored makes the platform exponentially more valuable. Winner-take-all dynamics.",
    "s4_feat4_title": "Data Licensing Ready", "s4_feat4_body": "Dish intelligence data ready to license to DoorDash, UberEats, Instacart, and CPG brands.",
    "s5_tam_label": "TAM", "s5_tam_val": "$371B", "s5_tam_desc": ["US restaurant industry", "and the race is on"],
    "s5_sam_label": "SAM", "s5_sam_val": "$23B", "s5_sam_desc": ["Restaurant tech market", "growing 12.5% CAGR"],
    "s5_som_label": "SOM", "s5_som_val": "$2.3B", "s5_som_desc": ["Dish intelligence market", "we are creating this category"],
    "s5_footer": "The window to own dish-level intelligence closes within 18 months",
    "s6_pct1": "40%", "s6_rev1": "B2C Subscriptions", "s6_rev1_desc": ["Viral-ready premium tier at $4.99/mo", "Winner-take-all consumer lock-in"],
    "s6_pct2": "35%", "s6_rev2": "B2B Restaurant SaaS", "s6_rev2_desc": ["$500/mo dish analytics dashboards", "Restaurants NEED this data"],
    "s6_pct3": "15%", "s6_rev3": "Data Licensing", "s6_rev3_desc": ["DoorDash, UberEats, Instacart all want this", "Enterprise data deals"],
    "s6_pct4": "10%", "s6_rev4": "Promoted Dishes", "s6_rev4_desc": ["High-intent, high-conversion ad placements", "CPC model"],
    "s7_desc": "No one is doing this YET. But Google, Yelp, and DoorDash all could. Our head start is everything.",
    "s8_p1_phase": "Phase 1", "s8_p1_name": "Ignite", "s8_p1_time": "Q2-Q4 2026",
    "s8_p1_items": ["Coordinated Philly influencer drops", "Founding user cohorts (limited spots)", "Restaurant FOMO campaign", "Press & media pipeline activated"],
    "s8_p2_phase": "Phase 2", "s8_p2_name": "Blitz", "s8_p2_time": "2027",
    "s8_p2_items": ["Tri-state expansion blitz", "8,000 premium users target", "Strategic restaurant chain partnerships", "Data licensing first deals"],
    "s8_p3_phase": "Phase 3", "s8_p3_name": "Dominate", "s8_p3_time": "2028-2029",
    "s8_p3_items": ["National rollout, 25 cities", "100K+ users, network effects locked in", "Category ownership cemented", "Acquisition conversations begin"],
    "s9_headline": "Already buzzing. Momentum is building.",
    "s9_s1_val": "2,490+", "s9_s1_label": "Dishes Scored & Live",
    "s9_s2_val": "9", "s9_s2_label": "Verticals Active",
    "s9_s3_val": "Live", "s9_s3_label": "App Store Approved",
    "s9_s4_val": "1", "s9_s4_label": "Patent Filed",
    "s9_s5_val": "93K", "s9_s5_label": "Lines of Code Shipped",
    "s10_m1_init": "NA", "s10_m1_name": "Nik Amin", "s10_m1_role": "CEO & Co-Founder",
    "s10_m1_desc": ["Bay Area engineer, built the entire platform", "Moving fast. 93K lines and counting."],
    "s10_m2_init": "AM", "s10_m2_name": "Ali Mabsoute", "s10_m2_role": "CMO & Co-Founder",
    "s10_m2_desc": ["UPenn, 15 years in Philly food scene", "Hyatt, Citi, AmEx. Knows the market cold."],
    "s10_m3_init": "", "s10_m3_name": "", "s10_m3_role": "", "s10_m3_desc": ["", ""],
    "s10_m4_init": "", "s10_m4_name": "", "s10_m4_role": "", "s10_m4_desc": ["", ""],
    "s10_advisory": "We chose Philadelphia. We chose to move fast. Founding investors get in at the ground floor.",
    "s11_v1": "$120K", "s11_l1": "2026 ARR",
    "s11_v2": "$1.2M", "s11_l2": "2027 ARR",
    "s11_v3": "$18M+", "s11_l3": "2029 ARR (Upside)",
    "s11_v4": "70%+", "s11_l4": "Gross Margin",
    "s11_v5": "3 months", "s11_l5": "CAC Payback",
    "s11_v6": "6.5x", "s11_l6": "LTV:CAC Ratio",
    "s11_footer": "Aggressive upside scenario  |  Winner-take-all economics amplify returns",
    "s12_amount": "$500K - $1.5M Seed",
    "s12_valuation": "at a $3M-$6M pre-money valuation",
    "s12_uof": "Use of Funds",
    "s12_pct1": "40%", "s12_uof1": "Engineering & AI", "s12_uof1d": "Widen the moat before big tech moves",
    "s12_pct2": "25%", "s12_uof2": "Growth & Marketing", "s12_uof2d": "Influencer drops, founding cohorts",
    "s12_pct3": "20%", "s12_uof3": "Business Development", "s12_uof3d": "Lock in restaurant & data partnerships",
    "s12_pct4": "15%", "s12_uof4": "Operations", "s12_uof4d": "Legal, patent prosecution",
    "s12_cta": "Our round is filling. 5 Founding Investor seats remain.  |  nik@tastyiq.com  |  tastyiq.com",
}

# ============================================================
# VARIATION 3: FOOD INDUSTRY B2B
# ============================================================
V3 = {
    "s1_brand": "Tastyr IQ",
    "s1_tagline": "The analytics layer restaurants never had.",
    "s1_subtitle": "AI-Powered Dish Intelligence for Restaurant Operators",
    "s1_meta": "Seed Round  |  March 2026  |  Confidential",
    "s2_title": "The Problem",
    "s2_headline": "Restaurants have no idea which dishes drive loyalty.",
    "s2_card1_title": "POS Data Is Not Dish Intelligence",
    "s2_card1_body": "Toast and Square tell you what sold. They cannot tell you what customers actually loved or hated.",
    "s2_card2_title": "Menu Optimization Is Guesswork",
    "s2_card2_body": "Chefs and operators make menu decisions based on intuition, not data. Margins suffer.",
    "s2_card3_title": "Review Noise, Zero Signal",
    "s2_card3_body": "Yelp and Google reviews mention dishes buried in paragraphs. No operator can extract actionable dish-level insights at scale.",
    "s2_stat": "78%",
    "s2_stat_desc": ["of restaurant operators say", "they lack actionable data", "on individual dish performance"],
    "s2_source": "Source: Restaurant Technology Survey 2025",
    "s3_desc": "Tastyr IQ gives restaurants an AI-powered analytics dashboard that scores every dish and reveals what drives customer loyalty.",
    "s3_f1_title": "Dish Performance Dashboard",
    "s3_f1_body": "Real-time scoring and trend analysis for every dish on the menu. See what customers love, what needs work, and what to cut.",
    "s3_f2_title": "Menu Optimization Engine",
    "s3_f2_body": "AI-powered recommendations for pricing, positioning, and seasonal rotations based on dish performance data.",
    "s3_f3_title": "Competitive Dish Benchmarking",
    "s3_f3_body": "See how your dishes score against competitors in your area. Identify gaps and opportunities.",
    "s3_f4_title": "Consumer Demand Signals",
    "s3_f4_body": "Understand what dishes consumers are searching for in your area before they walk through your door.",
    "s4_brand": "TASTYR IQ", "s4_search": "Search dishes, restaurants...", "s4_item_title": "Truffle Mushroom Risotto",
    "s4_dish1_sub": "Italian  \u2605 94/100", "s4_dish1_badge": "94",
    "s4_dish2_name": "Pan-Seared Salmon", "s4_dish2_sub": "Seafood  \u2605 92/100", "s4_dish2_badge": "92",
    "s4_dish3_name": "Wagyu Beef Burger", "s4_dish3_sub": "American  \u2605 89/100", "s4_dish3_badge": "89",
    "s4_feat1_title": "Dish Performance Dashboard", "s4_feat1_body": "Real-time scoring and trend analysis for every dish. See what customers love and what needs work.",
    "s4_feat2_title": "Menu Optimization Engine", "s4_feat2_body": "AI recommendations for pricing, positioning, and seasonal rotations based on dish data.",
    "s4_feat3_title": "Competitive Benchmarking", "s4_feat3_body": "See how your dishes score against competitors in your area. Identify gaps and opportunities.",
    "s4_feat4_title": "Consumer Demand Signals", "s4_feat4_body": "Understand what dishes consumers search for in your area before they walk through your door.",
    "s5_tam_label": "TAM", "s5_tam_val": "$23B", "s5_tam_desc": ["Restaurant technology", "and analytics SaaS market"],
    "s5_sam_label": "SAM", "s5_sam_val": "$5.8B", "s5_sam_desc": ["Restaurant analytics &", "menu intelligence tools"],
    "s5_som_label": "SOM", "s5_som_val": "$580M", "s5_som_desc": ["Dish-level analytics SaaS", "Year 5 target"],
    "s5_footer": "Restaurant SaaS market growing at 15.1% CAGR through 2030",
    "s6_pct1": "50%", "s6_rev1": "Restaurant Analytics SaaS", "s6_rev1_desc": ["Dish performance dashboards at $500/mo", "Tiered plans: Basic / Pro / Enterprise"],
    "s6_pct2": "25%", "s6_rev2": "Data Licensing", "s6_rev2_desc": ["Dish intelligence for delivery platforms & CPG", "Enterprise annual contracts"],
    "s6_pct3": "15%", "s6_rev3": "B2C Subscriptions", "s6_rev3_desc": ["Consumer app drives data collection", "$4.99/mo premium tier"],
    "s6_pct4": "10%", "s6_rev4": "POS Integration Fees", "s6_rev4_desc": ["Integration partnerships with Toast, Square", "Per-location setup fees"],
    "s7_desc": "Toast does POS. OpenTable does reservations. No one offers dish-level performance analytics. We are the first.",
    "s7_table": [
        ['Feature', 'Tastyr IQ', 'Toast', 'Yelp', 'OpenTable'],
        ['Dish-Level Scores', '\u2713', '\u2717', '\u2717', '\u2717'],
        ['AI Scoring Algorithm', '\u2713', '\u2717', '\u2717', '\u2717'],
        ['Menu Optimization', '\u2713', '\u2717', '\u2717', '\u2717'],
        ['Competitive Benchmarks', '\u2713', '\u2717', 'Partial', '\u2717'],
        ['Consumer Discovery', '\u2713', '\u2717', '\u2713', 'Partial'],
        ['POS Integration', 'Planned', '\u2713', '\u2717', '\u2717'],
        ['Data Licensing', '\u2713', '\u2717', '\u2717', '\u2717'],
    ],
    "s8_p1_phase": "Phase 1", "s8_p1_name": "Prove", "s8_p1_time": "Q2-Q4 2026",
    "s8_p1_items": ["Founding restaurant cohort: 10 Philly restaurants", "Prove dish analytics ROI with case studies", "POS integration pilot (Toast)", "Consumer app drives data flywheel"],
    "s8_p2_phase": "Phase 2", "s8_p2_name": "Partner", "s8_p2_time": "2027",
    "s8_p2_items": ["80 restaurant dashboards tri-state", "Distribution through POS partnerships", "Restaurant group enterprise deals", "First data licensing to delivery platforms"],
    "s8_p3_phase": "Phase 3", "s8_p3_name": "Platform", "s8_p3_time": "2028-2029",
    "s8_p3_items": ["800 restaurant dashboards nationally", "Embedded in POS workflows", "Enterprise data licensing at scale", "Category standard for dish analytics"],
    "s9_headline": "Product built. Restaurants are asking for this.",
    "s9_s1_val": "2,490+", "s9_s1_label": "Dishes Scored Across 9 Verticals",
    "s9_s2_val": "Live", "s9_s2_label": "Consumer App on App Store",
    "s9_s3_val": "93K", "s9_s3_label": "Lines of Production Code",
    "s9_s4_val": "Patent", "s9_s4_label": "Pending on Scoring Algorithm",
    "s9_s5_val": "$500/mo", "s9_s5_label": "Restaurant Dashboard Pricing",
    "s10_m1_init": "NA", "s10_m1_name": "Nik Amin", "s10_m1_role": "CEO & Co-Founder",
    "s10_m1_desc": ["Bay Area, Full-Stack Engineering", "Built entire platform end-to-end"],
    "s10_m2_init": "AM", "s10_m2_name": "Ali Mabsoute", "s10_m2_role": "CMO & Co-Founder",
    "s10_m2_desc": ["UPenn, Hyatt Hotels alumni", "15 years in Philadelphia restaurant scene"],
    "s10_m3_init": "", "s10_m3_name": "", "s10_m3_role": "", "s10_m3_desc": ["", ""],
    "s10_m4_init": "", "s10_m4_name": "", "s10_m4_role": "", "s10_m4_desc": ["", ""],
    "s10_advisory": "Advisors include restaurant operators, POS integration experts, and food industry data scientists.",
    "s11_v1": "$120K", "s11_l1": "2026 ARR",
    "s11_v2": "$1.2M", "s11_l2": "2027 ARR",
    "s11_v3": "$18M", "s11_l3": "2029 ARR (Target)",
    "s11_v4": "75%+", "s11_l4": "Gross Margin (SaaS)",
    "s11_v5": "4 months", "s11_l5": "CAC Payback (B2B)",
    "s11_v6": "8x", "s11_l6": "LTV:CAC Ratio",
    "s11_footer": "B2B SaaS unit economics  |  Restaurant dashboards are high-margin recurring revenue",
    "s12_amount": "$500K - $1.5M Seed",
    "s12_valuation": "at a $3M-$6M pre-money valuation",
    "s12_uof": "Use of Funds",
    "s12_pct1": "40%", "s12_uof1": "Engineering & AI", "s12_uof1d": "Restaurant dashboard, POS integrations",
    "s12_pct2": "25%", "s12_uof2": "Growth & Marketing", "s12_uof2d": "Restaurant sales, case studies, events",
    "s12_pct3": "20%", "s12_uof3": "Business Development", "s12_uof3d": "POS partnerships, chain restaurant deals",
    "s12_pct4": "15%", "s12_uof4": "Operations", "s12_uof4d": "Legal, IP, compliance",
    "s12_cta": "nik@tastyiq.com  |  tastyiq.com  |  The dish analytics platform restaurants have been waiting for.",
}

# ============================================================
# VARIATION 4: VISIONARY / DISRUPTOR
# ============================================================
V4 = {
    "s1_brand": "Tastyr IQ",
    "s1_tagline": "We're creating an entirely new data category.",
    "s1_subtitle": "Dish Intelligence: The Missing Layer of the Food Industry",
    "s1_meta": "Seed Round  |  March 2026  |  Confidential",
    "s2_title": "The Problem",
    "s2_headline": "For 20 years, the food industry has rated the wrong thing.",
    "s2_card1_title": "The Restaurant Fallacy",
    "s2_card1_body": "We rate restaurants. But people eat dishes. The entire data model of the food industry is built on the wrong unit of analysis.",
    "s2_card2_title": "No Data Layer Exists",
    "s2_card2_body": "Movies have Rotten Tomatoes. Music has Spotify. Hotels have star ratings per room type. Food has nothing at the dish level.",
    "s2_card3_title": "A $371B Industry, Data-Blind",
    "s2_card3_body": "The largest consumer industry in America has zero standardized data infrastructure for its core product: the dish.",
    "s2_stat": "20 yrs",
    "s2_stat_desc": ["of rating restaurants", "when we should have been", "rating dishes"],
    "s2_source": "The paradigm is about to shift.",
    "s3_desc": "Tastyr IQ is creating dish intelligence -- an entirely new data category that will become infrastructure for the food industry.",
    "s3_f1_title": "A New Data Category",
    "s3_f1_body": "Dish intelligence does not exist today. We are building the foundational data layer that the entire food ecosystem will depend on.",
    "s3_f2_title": "Patent-Pending AI",
    "s3_f2_body": "Our proprietary scoring algorithm is patent-pending. Category ownership through intellectual property, not just execution speed.",
    "s3_f3_title": "Three-Sided Marketplace",
    "s3_f3_body": "Consumers discover. Restaurants optimize. Data partners license. Three value streams from one intelligence platform.",
    "s3_f4_title": "9 Verticals Deep",
    "s3_f4_body": "Not a horizontal play. Deep vertical expertise across 9 cuisines, understanding what excellence means in each category.",
    "s4_brand": "TASTYR IQ", "s4_search": "Discover the best dishes...", "s4_item_title": "Wagyu Beef Tartare",
    "s4_dish1_sub": "Japanese  \u2605 98/100", "s4_dish1_badge": "98",
    "s4_dish2_name": "Lamb Shank Tagine", "s4_dish2_sub": "Moroccan  \u2605 95/100", "s4_dish2_badge": "95",
    "s4_dish3_name": "Duck Confit", "s4_dish3_sub": "French  \u2605 93/100", "s4_dish3_badge": "93",
    "s4_feat1_title": "A New Data Category", "s4_feat1_body": "Dish intelligence doesn't exist today. We're building the foundational data layer the food ecosystem will depend on.",
    "s4_feat2_title": "Patent-Pending AI", "s4_feat2_body": "Category ownership through intellectual property. Our scoring algorithm is patent-pending.",
    "s4_feat3_title": "Three-Sided Marketplace", "s4_feat3_body": "Consumers discover. Restaurants optimize. Data partners license. Three value streams from one platform.",
    "s4_feat4_title": "9 Verticals Deep", "s4_feat4_body": "Deep vertical expertise across 9 cuisines, understanding what excellence means in each category.",
    "s5_tam_label": "TAM", "s5_tam_val": "$371B", "s5_tam_desc": ["US restaurant industry", "zero dish-level data today"],
    "s5_sam_label": "SAM", "s5_sam_val": "$23B", "s5_sam_desc": ["Restaurant tech market", "no dish intelligence player"],
    "s5_som_label": "SOM", "s5_som_val": "$2.3B", "s5_som_desc": ["New category: dish intelligence", "we define the market"],
    "s5_footer": "We are not entering a market. We are creating one.",
    "s6_pct1": "35%", "s6_rev1": "Consumer Platform", "s6_rev1_desc": ["Premium dish discovery at $4.99/mo", "Category-defining consumer product"],
    "s6_pct2": "35%", "s6_rev2": "Restaurant Intelligence", "s6_rev2_desc": ["$500/mo dish analytics SaaS", "Menu optimization & benchmarking"],
    "s6_pct3": "20%", "s6_rev3": "Data Infrastructure", "s6_rev3_desc": ["License dish intelligence to the food ecosystem", "Delivery, CPG, hospitality, media"],
    "s6_pct4": "10%", "s6_rev4": "Platform Services", "s6_rev4_desc": ["Promoted dishes, integrations", "Ecosystem revenue"],
    "s7_desc": "There is no competition because this category does not exist yet. We are building it.",
    "s8_p1_phase": "Phase 1", "s8_p1_name": "Prove", "s8_p1_time": "Q2-Q4 2026",
    "s8_p1_items": ["Philadelphia is our laboratory", "Prove dish intelligence has value", "1,000 founding users", "10 restaurant data partnerships"],
    "s8_p2_phase": "Phase 2", "s8_p2_name": "Build", "s8_p2_time": "2027",
    "s8_p2_items": ["Tri-state expansion", "Category awareness campaign", "First data licensing partnerships", "Industry recognition and press"],
    "s8_p3_phase": "Phase 3", "s8_p3_name": "Lead", "s8_p3_time": "2028-2029",
    "s8_p3_items": ["National: the nation is our market", "Category leadership cemented", "Dish intelligence becomes industry standard", "Strategic acquisition interest"],
    "s9_headline": "Category creation in progress.",
    "s9_s1_val": "2,490+", "s9_s1_label": "Dishes in Intelligence Layer",
    "s9_s2_val": "9", "s9_s2_label": "Cuisine Verticals Deep",
    "s9_s3_val": "Patent", "s9_s3_label": "Pending = Category Ownership",
    "s9_s4_val": "93K", "s9_s4_label": "Lines of Code = Serious Engineering",
    "s9_s5_val": "Live", "s9_s5_label": "App Store + TestFlight Beta",
    "s10_m1_init": "NA", "s10_m1_name": "Nik Amin", "s10_m1_role": "CEO & Co-Founder",
    "s10_m1_desc": ["Bay Area engineer who chose to build", "93K lines of code, patent-pending AI"],
    "s10_m2_init": "AM", "s10_m2_name": "Ali Mabsoute", "s10_m2_role": "CMO & Co-Founder",
    "s10_m2_desc": ["UPenn grad who chose Philadelphia", "Hyatt, Citi, AmEx -- 15yr market expertise"],
    "s10_m3_init": "", "s10_m3_name": "", "s10_m3_role": "", "s10_m3_desc": ["", ""],
    "s10_m4_init": "", "s10_m4_name": "", "s10_m4_role": "", "s10_m4_desc": ["", ""],
    "s10_advisory": "A UPenn-adjacent grad and a Bay Area engineer chose to build dish intelligence from Philadelphia. This is a founding story.",
    "s11_v1": "$120K", "s11_l1": "2026 ARR",
    "s11_v2": "$1.2M", "s11_l2": "2027 ARR",
    "s11_v3": "$18M", "s11_l3": "2029 ARR (Target)",
    "s11_v4": "70%+", "s11_l4": "Gross Margin",
    "s11_v5": "3 months", "s11_l5": "CAC Payback",
    "s11_v6": "6.5x", "s11_l6": "LTV:CAC Ratio",
    "s11_footer": "The path to category leadership  |  First-mover economics in a new data category",
    "s12_amount": "$500K - $1.5M Seed",
    "s12_valuation": "at a $3M-$6M pre-money valuation",
    "s12_uof": "Use of Funds",
    "s12_pct1": "40%", "s12_uof1": "Engineering & AI", "s12_uof1d": "Deepen the intelligence layer",
    "s12_pct2": "25%", "s12_uof2": "Growth & Marketing", "s12_uof2d": "Category awareness, founding users",
    "s12_pct3": "20%", "s12_uof3": "Business Development", "s12_uof3d": "Data partnerships, restaurant onboarding",
    "s12_pct4": "15%", "s12_uof4": "Operations", "s12_uof4d": "Patent prosecution, legal, admin",
    "s12_cta": "Join us at the ground floor of dish intelligence. This is a once-in-a-generation category creation moment.  |  nik@tastyiq.com",
}

# ============================================================
# BUILD ALL 4
# ============================================================
decks = [
    ("pitch-deck-v1-investor.pptx", V1),
    ("pitch-deck-v2-fomo.pptx", V2),
    ("pitch-deck-v3-b2b.pptx", V3),
    ("pitch-deck-v4-visionary.pptx", V4),
]

for filename, v in decks:
    print(f"Building: {filename}...")
    build_deck(filename, v)

print("\n=== ALL 4 PITCH DECKS BUILT SUCCESSFULLY ===")
