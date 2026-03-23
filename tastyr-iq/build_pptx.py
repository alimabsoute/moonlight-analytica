#!/usr/bin/env python3
"""
Tastyr IQ — PPTX Deck Builder
Creates a 20-slide professional presentation deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import csv
import os

NAVY = RGBColor(0x0A, 0x1A, 0x2F)
GOLD = RGBColor(0xD4, 0xA5, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_NAVY = RGBColor(0x1A, 0x2A, 0x44)
ACCENT_BLUE = RGBColor(0x3A, 0x7C, 0xD5)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=12, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return txBox

def add_gold_bar(slide, top):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top), Inches(0.08), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: TITLE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)
    # Gold accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.333), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    add_text_box(slide, 1, 1.5, 11, 1.5, "TASTYR IQ", 54, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.6, 11, 0.8, "Tri-State Outreach & Investment Strategy", 28, WHITE, False, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 5, 11, 0.5, "AI-Powered Dish Intelligence  |  Patent Pending  |  App Store Live", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 6, 11, 0.5, "March 2026  |  Confidential", 14, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

    # ========== SLIDE 2: THE OPPORTUNITY ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 5, 0.7, "THE OPPORTUNITY", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    add_text_box(slide, 0.5, 1.3, 6, 0.6, "$371B US Restaurant Industry", 24, WHITE, True)
    add_bullet_list(slide, 0.5, 2, 6, 4, [
        "Nobody rates the actual FOOD — only restaurants",
        "A 4-star restaurant can have 2-star dishes",
        "Consumers waste money on bad dishes at good restaurants",
        "Restaurants don't know which dishes are actually loved",
        "No data layer exists at the dish level",
        "",
        "Tastyr IQ is Rotten Tomatoes for food.",
    ], 16, WHITE)

    # Right side - market size boxes
    for i, (label, val) in enumerate([("TAM", "$371B"), ("SAM", "$23B"), ("SOM", "$2.3B")]):
        y = 1.5 + i * 1.6
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(y), Inches(4.5), Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = GOLD
        shape.line.width = Pt(1)
        add_text_box(slide, 8.3, y + 0.1, 4, 0.4, label, 14, GOLD, True)
        add_text_box(slide, 8.3, y + 0.5, 4, 0.6, val, 36, WHITE, True)

    # ========== SLIDE 3: OUR TRACTION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 5, 0.7, "OUR TRACTION", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    metrics = [
        ("9", "Live Dish\nVerticals"),
        ("2,490+", "Indexed\nDishes"),
        ("~93K", "Lines of\nCode"),
        ("Patent", "Pending\nFiling"),
        ("Live", "App Store\nApproved"),
    ]
    for i, (val, label) in enumerate(metrics):
        x = 0.5 + i * 2.5
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(2.2), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = GOLD
        shape.line.width = Pt(1)
        add_text_box(slide, x + 0.1, 1.7, 2, 1, val, 40, GOLD, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.1, 3, 2, 0.8, label, 14, WHITE, False, PP_ALIGN.CENTER)

    add_text_box(slide, 0.5, 4.5, 12, 1.5,
        "Co-Founders: Nik Amin (CEO, Bay Area) & Ali Mabsoute (CMO, Philadelphia)\n"
        "Ali: UPenn | 15yr Philadelphia | Hyatt, Citibank, American Express | Old City Labs",
        14, LIGHT_GRAY, False, PP_ALIGN.LEFT)

    # ========== SLIDE 4: OUTREACH PHILOSOPHY ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "OUTREACH PHILOSOPHY", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    add_text_box(slide, 1, 1.5, 11, 1, '"Create the wave. Don\'t chase it."', 36, WHITE, True, PP_ALIGN.CENTER)

    pillars = [
        ("SOCIAL PROOF", "Build buzz before\napproaching investors.\nInfluencers, media,\nrestaurateurs first."),
        ("CHAIN STRATEGY", "Every person is a link.\nBlogger > Restaurant >\nInvestor. Never cold."),
        ("PSYCHOLOGY", "Cialdini principles\napplied to every\nmessage. Authority,\nScarcity, Unity."),
    ]
    for i, (title, desc) in enumerate(pillars):
        x = 0.8 + i * 4.2
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3), Inches(3.8), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = GOLD
        shape.line.width = Pt(1.5)
        add_text_box(slide, x + 0.2, 3.2, 3.4, 0.5, title, 18, GOLD, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 3.9, 3.4, 2.5, desc, 14, WHITE, False, PP_ALIGN.CENTER)

    # ========== SLIDE 5: THE 12 SEGMENTS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "THE 12 SEGMENTS", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    # Read CRM for counts
    crm_path = os.path.join(os.path.dirname(__file__), "crm", "master-crm.csv")
    segment_counts = {}
    with open(crm_path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            seg = row["Segment"]
            segment_counts[seg] = segment_counts.get(seg, 0) + 1

    segments = [
        ("INVESTOR TRACK", [
            ("Angel Investors", segment_counts.get("Angel", 0)),
            ("Venture Capital", segment_counts.get("VC", 0)),
            ("Private Equity", segment_counts.get("PE", 0)),
            ("Accelerators", segment_counts.get("Accelerator", 0)),
        ]),
        ("INFLUENCE TRACK", [
            ("Restaurateurs", segment_counts.get("Restaurateur", 0)),
            ("Food Bloggers", segment_counts.get("Food Blogger", 0)),
            ("Media", segment_counts.get("Media", 0)),
            ("University", segment_counts.get("University", 0)),
        ]),
        ("CONNECTOR TRACK", [
            ("Tech Founders", segment_counts.get("Founder-Connector", 0)),
            ("Food Industry", segment_counts.get("Food Industry", 0)),
            ("NJ/DE/Suburbs", segment_counts.get("NJ-DE-Suburb", 0)),
        ]),
    ]

    col_x = [0.5, 4.7, 8.9]
    for col, (track_name, segs) in enumerate(segments):
        x = col_x[col]
        add_text_box(slide, x, 1.3, 4, 0.5, track_name, 16, GOLD, True)
        for i, (name, count) in enumerate(segs):
            y = 2 + i * 1.2
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(1))
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_NAVY
            shape.line.color.rgb = GOLD if count > 25 else ACCENT_BLUE
            shape.line.width = Pt(1)
            add_text_box(slide, x + 0.2, y + 0.1, 2.5, 0.4, name, 14, WHITE, True)
            add_text_box(slide, x + 0.2, y + 0.5, 2.5, 0.4, f"{count} contacts", 12, GOLD)

    total = sum(segment_counts.values())
    add_text_box(slide, 0.5, 6.5, 12, 0.5, f"TOTAL: {total} contacts across the tri-state area (PA, NJ, DE)", 16, GOLD, True, PP_ALIGN.CENTER)

    # ========== SLIDE 6: GEOGRAPHIC COVERAGE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "GEOGRAPHIC COVERAGE", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    # Count by state
    state_counts = {"PA": 0, "NJ": 0, "DE": 0, "NY": 0, "Other": 0}
    with open(crm_path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            loc = row.get("City/State", "").lower()
            if "pa" in loc or "philadelphia" in loc or "philly" in loc or "pennsylvania" in loc:
                state_counts["PA"] += 1
            elif "nj" in loc or "new jersey" in loc or "jersey" in loc or "princeton" in loc or "hoboken" in loc or "newark" in loc:
                state_counts["NJ"] += 1
            elif "de" in loc or "delaware" in loc or "wilmington" in loc:
                state_counts["DE"] += 1
            elif "ny" in loc or "new york" in loc or "nyc" in loc:
                state_counts["NY"] += 1
            else:
                state_counts["Other"] += 1

    regions = [
        ("PHILADELPHIA (PA)", state_counts["PA"], "Primary Market\nAli's home base\n15 years of relationships", 0.5, 1.5),
        ("NEW JERSEY", state_counts["NJ"], "Secondary Market\nPrinceton corridor\nNJ Angels + food scene", 4.7, 1.5),
        ("DELAWARE", state_counts["DE"], "Tertiary Market\nHorn Entrepreneurship\nDogfish Head crossover", 8.9, 1.5),
        ("NYC SPILLOVER", state_counts["NY"], "Strategic Targets\nFood-tech VCs\nCollaborative, S2G, Lerer", 0.5, 4.5),
        ("SUBURBS & OTHER", state_counts["Other"], "Supporting Markets\nMain Line, Bucks County\nLehigh Valley", 4.7, 4.5),
    ]

    for name, count, desc, x, y in regions:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = GOLD
        shape.line.width = Pt(1)
        add_text_box(slide, x + 0.2, y + 0.15, 3.4, 0.4, name, 16, GOLD, True)
        add_text_box(slide, x + 0.2, y + 0.6, 3.4, 0.6, f"{count} contacts", 28, WHITE, True)
        add_text_box(slide, x + 0.2, y + 1.3, 3.4, 1, desc, 11, LIGHT_GRAY)

    # ========== SLIDE 7: RANKING SYSTEM ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "THE RANKING SYSTEM (1-100)", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    add_text_box(slide, 0.5, 1.3, 12, 0.6,
        "RANK = (Investor Potential x 0.30) + (Network Reach x 0.25) + (FOMO Multiplier x 0.20) + (Accessibility x 0.15) + (Timing x 0.10)",
        13, ACCENT_BLUE, False, PP_ALIGN.CENTER)

    tiers = [
        ("90-100", "Contact THIS WEEK", "148 contacts", "Drop everything. High-probability, high-impact."),
        ("70-89", "Contact within 2 weeks", "", "Strong fits, clear path to value."),
        ("50-69", "Contact within month 1", "150 contacts", "Good fits, need some groundwork."),
        ("30-49", "Month 2-3", "", "Build relationship first."),
        ("1-29", "Long game", "11 contacts", "Track, engage on social, wait."),
    ]

    # Count actual tiers
    tier_counts = {"90-100": 0, "70-89": 0, "50-69": 0, "30-49": 0, "1-29": 0}
    with open(crm_path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            r = int(row["Rank"])
            if r >= 90: tier_counts["90-100"] += 1
            elif r >= 70: tier_counts["70-89"] += 1
            elif r >= 50: tier_counts["50-69"] += 1
            elif r >= 30: tier_counts["30-49"] += 1
            else: tier_counts["1-29"] += 1

    for i, (rng, action, _, desc) in enumerate(tiers):
        y = 2.2 + i * 1
        count = tier_counts[rng]
        # Color gradient from gold to blue
        bar_color = GOLD if i < 2 else ACCENT_BLUE if i < 4 else MEDIUM_GRAY
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = bar_color
        shape.line.width = Pt(1)
        add_text_box(slide, 0.7, y + 0.1, 1.5, 0.5, rng, 18, bar_color, True)
        add_text_box(slide, 2.5, y + 0.1, 3, 0.5, action, 14, WHITE, True)
        add_text_box(slide, 6, y + 0.1, 2, 0.5, f"{count} contacts", 14, GOLD)
        add_text_box(slide, 8.5, y + 0.1, 4, 0.5, desc, 11, LIGHT_GRAY)

    # ========== SLIDES 8-11: CHAIN STRATEGIES ==========
    chains = [
        ("CHAIN: RESTAURATEUR GATEWAY", [
            "Step 1: Give food blogger exclusive beta access",
            "Step 2: Blogger posts dish scores at popular restaurant",
            "Step 3: Restaurant owner sees coverage, gets curious",
            "Step 4: Ali reaches out: 'Your khao soi scored 94/100'",
            "Step 5: Restaurateur becomes champion, tells investor friends",
            "Step 6: Investor hears from TRUSTED source, takes Ali's call",
        ], "Blogger > Restaurant > Investor\nBuilds organic social proof that pulls investors in."),
        ("CHAIN: FOMO CASCADE", [
            "Step 1: Seed 5-10 influencers with exclusive beta access",
            "Step 2: Coordinated posts in 48-hour window",
            "Step 3: 'Everyone's talking about this' perception created",
            "Step 4: Media picks up the trend (Eater, PhillyMag)",
            "Step 5: Restaurateurs scramble to get dishes rated",
            "Step 6: Investors see MARKET PULL, not push",
        ], "Perception of unstoppable momentum.\nInvestors arrive pre-warmed and excited."),
        ("CHAIN: UNIVERSITY PIPELINE", [
            "Step 1: Connect with Penn Food Innovation Lab faculty",
            "Step 2: Guest lecture / demo for students",
            "Step 3: Students become beta testers (hundreds of food-obsessed kids)",
            "Step 4: Usage data spike + campus buzz",
            "Step 5: Faculty intro to Wharton entrepreneurship professors",
            "Step 6: Penn alumni angel network invitation",
        ], "Academic credibility > Student traction > Investor intros.\nAli's UPenn connection is the entry point."),
        ("CHAIN: NJ-TO-NYC BRIDGE", [
            "Step 1: Connect with Princeton Entrepreneurship Council",
            "Step 2: Demo at monthly meetup or pitch competition",
            "Step 3: Princeton-area angels show interest",
            "Step 4: NJ angels co-invest with NYC-adjacent funds",
            "Step 5: NYC VC sees tri-state traction, not just 'a Philly app'",
            "Step 6: Collaborative Fund, Lerer Hippeau warm intro via NJ bridge",
        ], "Use NJ ecosystem to access NYC capital.\nShow tri-state traction, not just Philadelphia."),
    ]

    for title, steps, summary in chains:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, NAVY)
        add_text_box(slide, 0.5, 0.3, 10, 0.7, title, 28, GOLD, True)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
        shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

        # Steps as connected boxes
        for i, step in enumerate(steps):
            y = 1.3 + i * 0.85
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(8), Inches(0.7))
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_NAVY
            shape.line.color.rgb = GOLD if i < 3 else ACCENT_BLUE
            shape.line.width = Pt(1)
            add_text_box(slide, 0.7, y + 0.1, 7.6, 0.5, step, 13, WHITE)
            # Arrow connector
            if i < len(steps) - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.5), Inches(y + 0.7), Inches(0.3), Inches(0.15))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = GOLD
                arrow.line.fill.background()

        add_text_box(slide, 9, 1.5, 3.8, 3, summary, 14, LIGHT_GRAY)

    # ========== SLIDE 12: FOMO PLAYBOOK ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "FOMO PLAYBOOK OVERVIEW", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    tactics = [
        ("SCARCITY", ['"Only 50 spots in founding cohort"', '"Founding Investor — only 5 seats"', '"We\'re meeting with 10 investors max"', '"Round closing by [date]"']),
        ("SOCIAL PROOF", ['"Coordinated influencer drop (48hr window)"', '"As Seen In press momentum"', '"Restaurateur case studies"', '"App Store presence as credibility"']),
        ("EXCLUSIVITY", ['"Invite-only beta access"', '"Pre-patent NDA mystique"', '"Chosen / specifically selected language"', '"Embargo offers to media"']),
    ]

    for i, (title, items) in enumerate(tactics):
        x = 0.5 + i * 4.2
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(3.8), Inches(5.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = GOLD
        shape.line.width = Pt(1.5)
        add_text_box(slide, x + 0.2, 1.5, 3.4, 0.5, title, 20, GOLD, True, PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_text_box(slide, x + 0.3, 2.3 + j * 1.1, 3.2, 0.9, item, 12, WHITE)

    # ========== SLIDE 13: PSYCHOLOGY-DRIVEN OUTREACH ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "PSYCHOLOGY-DRIVEN OUTREACH", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    add_text_box(slide, 0.5, 1.2, 12, 0.5, "Every message engineered using Cialdini's 6 Principles of Influence", 16, WHITE, False, PP_ALIGN.CENTER)

    principles = [
        ("AUTHORITY", "Lead with UPenn, patent,\nHyatt/Citi/AmEx credentials"),
        ("SOCIAL PROOF", '"Other investors are\nalready in conversations"'),
        ("UNITY", '"Fellow Philadelphian\nbuilding for our city"'),
        ("SCARCITY", '"Limited beta / limited\ninvestor seats"'),
        ("RECIPROCITY", "Offer dish data/insights\nbefore asking for anything"),
        ("COMMITMENT", '"You\'ve said you want to\nback X — here\'s X"'),
    ]

    for i, (principle, desc) in enumerate(principles):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 2 + row * 2.5
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = GOLD
        shape.line.width = Pt(1)
        add_text_box(slide, x + 0.2, y + 0.2, 3.4, 0.4, principle, 16, GOLD, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, y + 0.8, 3.4, 1.2, desc, 13, WHITE, False, PP_ALIGN.CENTER)

    # ========== SLIDE 14: TIER 1 PRIORITY LIST ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "TIER 1 PRIORITY CONTACTS", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    # Read top 15 from CRM
    with open(crm_path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        all_rows = sorted(list(reader), key=lambda r: int(r["Rank"]), reverse=True)
        top15 = all_rows[:15]

    # Table header
    add_text_box(slide, 0.5, 1.2, 1, 0.4, "RANK", 10, GOLD, True)
    add_text_box(slide, 1.5, 1.2, 3, 0.4, "NAME", 10, GOLD, True)
    add_text_box(slide, 4.5, 1.2, 3, 0.4, "ORGANIZATION", 10, GOLD, True)
    add_text_box(slide, 7.5, 1.2, 2, 0.4, "SEGMENT", 10, GOLD, True)
    add_text_box(slide, 9.5, 1.2, 3, 0.4, "CHAIN", 10, GOLD, True)

    for i, row in enumerate(top15):
        y = 1.6 + i * 0.35
        bg_color = LIGHT_NAVY if i % 2 == 0 else NAVY
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(12), Inches(0.33))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        add_text_box(slide, 0.5, y, 1, 0.33, row["Rank"], 10, GOLD, True)
        add_text_box(slide, 1.5, y, 3, 0.33, row["Name"][:30], 10, WHITE)
        add_text_box(slide, 4.5, y, 3, 0.33, row["Organization"][:35], 10, LIGHT_GRAY)
        add_text_box(slide, 7.5, y, 2, 0.33, row["Segment"], 10, ACCENT_BLUE)
        add_text_box(slide, 9.5, y, 3, 0.33, row["Chain Strategy"][:30], 10, LIGHT_GRAY)

    # ========== SLIDE 15: WEEK-BY-WEEK ACTION PLAN ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "WEEK-BY-WEEK ACTION PLAN", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    weeks = [
        ("WEEK 1-2: FOUNDATION", [
            "Seed 5-10 influencers with beta access",
            "Onboard 3-5 champion restaurateurs",
            "Submit accelerator applications (LIFT Labs, Dreamit, Food-X)",
            "Begin Tier 1 outreach (Rank 90-100)",
            "Secure first media contact",
        ]),
        ("WEEK 3-4: MOMENTUM", [
            "Coordinated influencer drop (48-hour window)",
            "Media coverage from Eater Philly or Technical.ly",
            "University partnerships initiated",
            "Tier 2 outreach begins (Rank 70-89)",
            "NJ-to-NYC bridge chain activated",
        ]),
        ("WEEK 5-8: ACCELERATION", [
            "Multiple media mentions creating buzz",
            "Restaurateur case studies documented",
            "Student beta tester cohort active",
            "Investor meetings and term sheet discussions",
            "Full FOMO messaging: 'Round is filling'",
        ]),
        ("WEEK 9-12: CLOSE", [
            "Term sheet(s) received",
            "Strategic investors brought in",
            "Round closed at $500K-$1.5M",
            "Press: 'Tastyr IQ raises seed round'",
            "Prepare for Series A path",
        ]),
    ]

    for i, (title, items) in enumerate(weeks):
        x = 0.3 + i * 3.25
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(3), Inches(5.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = GOLD if i < 2 else ACCENT_BLUE
        shape.line.width = Pt(1)
        add_text_box(slide, x + 0.15, 1.5, 2.7, 0.6, title, 12, GOLD, True, PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_text_box(slide, x + 0.2, 2.3 + j * 0.9, 2.6, 0.8, f"  {item}", 10, WHITE)

    # ========== SLIDE 16: KPIs & SUCCESS METRICS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "KPIs & SUCCESS METRICS", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    kpis = [
        ("Investor Meetings", "15+", "within first 4 weeks"),
        ("Term Sheets", "2-3", "by week 8-10"),
        ("Influencer Posts", "10+", "in first 2 weeks"),
        ("Media Mentions", "5+", "across tri-state press"),
        ("Restaurateur Champions", "10+", "in founding cohort"),
        ("University Partners", "2-3", "active partnerships"),
        ("App Downloads", "500+", "in first month"),
        ("Round Closed", "$500K-$1.5M", "by week 12"),
    ]

    for i, (metric, target, timeline) in enumerate(kpis):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.5
        y = 1.3 + row * 1.4
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = GOLD
        shape.line.width = Pt(1)
        add_text_box(slide, x + 0.2, y + 0.1, 3, 0.4, metric, 14, WHITE, True)
        add_text_box(slide, x + 3.5, y + 0.1, 2, 0.4, target, 20, GOLD, True)
        add_text_box(slide, x + 0.2, y + 0.6, 5.5, 0.4, timeline, 11, LIGHT_GRAY)

    # ========== SLIDE 17: CRM OVERVIEW ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "CRM OVERVIEW", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    add_text_box(slide, 0.5, 1.3, 12, 0.5, f"{total} contacts | 24 columns | 12 segments | Google Sheets ready", 18, WHITE, True, PP_ALIGN.CENTER)

    blocks = [
        ("IDENTITY BLOCK", "Segment, Name, Title/Role,\nOrganization, City/State"),
        ("CONTACT BLOCK", "Email, LinkedIn,\nOther Social, Phone"),
        ("RELEVANCE BLOCK", "Investment Focus, Check Size,\nNotable Work, Food Connection"),
        ("PSYCHOLOGY BLOCK", "Why Contact, Psychological\nAngle, Emotional Trigger"),
        ("STRATEGY BLOCK", "Connection Degree, Warm Intro,\nRank, Chain Strategy, Approach,\nCustom Message, Status, Notes"),
    ]

    for i, (title, desc) in enumerate(blocks):
        x = 0.3 + (i % 3) * 4.3
        y = 2.2 if i < 3 else 4.5
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_NAVY
        shape.line.color.rgb = GOLD
        shape.line.width = Pt(1)
        add_text_box(slide, x + 0.2, y + 0.15, 3.6, 0.4, title, 14, GOLD, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, y + 0.7, 3.6, 1, desc, 12, WHITE, False, PP_ALIGN.CENTER)

    # ========== SLIDE 18: NEXT STEPS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "NEXT STEPS", 32, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    ali_tasks = [
        "Begin Tier 1 outreach this week (Rank 90-100 contacts)",
        "Send beta access invites to top 10 food influencers",
        "Schedule 3-5 restaurateur meetings",
        "Submit Comcast LIFT Labs and Food-X accelerator applications",
        "Pitch Eater Philly and Technical.ly for launch coverage",
        "Connect with Penn Food Innovation Lab faculty",
        "LinkedIn content: 'Building something exciting in Philly food-tech'",
    ]

    nik_tasks = [
        "Prepare investor demo environment",
        "Finalize one-pager and pitch deck for investor meetings",
        "Set up analytics dashboard for tracking user engagement",
        "Prepare restaurateur dashboard mockup for founding cohort",
        "Review and refine chain strategy timing with Ali",
    ]

    # Ali column
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_NAVY; shape.line.color.rgb = GOLD; shape.line.width = Pt(1)
    add_text_box(slide, 0.7, 1.5, 5.6, 0.5, "ALI (CMO — Philadelphia)", 18, GOLD, True)
    for i, task in enumerate(ali_tasks):
        add_text_box(slide, 0.8, 2.2 + i * 0.65, 5.4, 0.6, f"  {task}", 12, WHITE)

    # Nik column
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.3), Inches(5.8), Inches(5.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_NAVY; shape.line.color.rgb = ACCENT_BLUE; shape.line.width = Pt(1)
    add_text_box(slide, 7.2, 1.5, 5.4, 0.5, "NIK (CEO — Bay Area)", 18, ACCENT_BLUE, True)
    for i, task in enumerate(nik_tasks):
        add_text_box(slide, 7.3, 2.2 + i * 0.65, 5.2, 0.6, f"  {task}", 12, WHITE)

    # ========== SLIDE 19: SEGMENT BREAKDOWN ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_text_box(slide, 0.5, 0.3, 10, 0.7, "APPENDIX: SEGMENT BREAKDOWN", 28, GOLD, True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(5), Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    seg_details = [
        ("Angel", "Robin Hood Ventures, Broad St Angels, Keiretsu, NJ Angels, Princeton Angel Group, Golden Seeds, Delaware Crossing"),
        ("VC", "First Round, Dreamit, Osage, SeventySix, Collaborative Fund, S2G Ventures, Lerer Hippeau, Big Idea Ventures"),
        ("PE", "Graham Partners, Swander Pace, Mistral Equity, L Catterton, Brynwood Partners"),
        ("Accelerator", "Comcast LIFT Labs, Big Idea Ventures, EatOkra, Plug and Play, Horn Entrepreneurship"),
        ("Restaurateur", "Vetri, Solomonov/CookNSolo, Starr, Yin, Garces, Vernick, Burke, Calagione + 24 more"),
        ("Food Blogger", "PhillyFoodFanatic (577K), JoshEatsPhilly (190K), PhillyFoodies (135K) + 28 more"),
        ("Media", "Craig LaBan (Inquirer), Kae Lani Palmisano (PhillyMag), Technical.ly + 22 more"),
        ("University", "UPenn/Wharton, Temple, Drexel, Princeton, Rutgers, U of Delaware"),
        ("Founder-Connector", "PSL, 1Philadelphia, DuckDuckGo, Crossbeam, Guru, Black and Mobile + 18 more"),
        ("Food Industry", "Sysco, US Foods, Baldor, Toast, DoorDash, PRLA, NJRHA, DRA"),
        ("NJ-DE-Suburb", "South Jersey, Princeton corridor, Wilmington, Main Line, Bucks County, Lehigh Valley"),
    ]

    for i, (seg, orgs) in enumerate(seg_details):
        y = 1.3 + i * 0.52
        count = segment_counts.get(seg, 0)
        add_text_box(slide, 0.5, y, 2, 0.5, f"{seg}", 11, GOLD, True)
        add_text_box(slide, 2.5, y, 1, 0.5, f"{count}", 11, WHITE, True)
        add_text_box(slide, 3.5, y, 9, 0.5, orgs[:100], 9, LIGHT_GRAY)

    # ========== SLIDE 20: CLOSING ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), Inches(13.333), Inches(0.04))
    shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

    add_text_box(slide, 1, 1.5, 11, 1, "TASTYR IQ", 54, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.7, 11, 0.8, "Rates Dishes, Not Restaurants.", 28, WHITE, False, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4, 11, 0.6, "ali@tastyriq.com  |  tastyriq.com", 18, ACCENT_BLUE, False, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 5, 11, 0.5, "Nik Amin (CEO)  &  Ali Mabsoute (CMO)", 16, WHITE, False, PP_ALIGN.CENTER)
    add_text_box(slide, 1, 5.8, 11, 0.5, "Confidential — March 2026", 14, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "strategy", "tastyr-iq-outreach-deck.pptx")
    prs.save(output_path)
    print(f"PPTX saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
